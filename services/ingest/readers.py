"""Hexis ingestion — split from the former services/ingest.py (#89).
Module: readers.
"""

from __future__ import annotations

import argparse
import asyncio
import hashlib
import json
import os
import re
import sys
import time
from dataclasses import dataclass, field
from datetime import datetime, timezone
from enum import Enum
from pathlib import Path
from typing import Any, Callable, Optional
from uuid import UUID

from .config import Section, _extract_title

# =========================================================================
# DOCUMENT READERS
# =========================================================================


class DocumentReader:
    @staticmethod
    def read(file_path: Path) -> str:
        raise NotImplementedError


class MarkdownReader(DocumentReader):
    @staticmethod
    def read(file_path: Path) -> str:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()


class TextReader(DocumentReader):
    @staticmethod
    def read(file_path: Path) -> str:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()


class CodeReader(DocumentReader):
    LANGUAGE_MAP = {
        ".py": "python",
        ".js": "javascript",
        ".ts": "typescript",
        ".jsx": "javascript-react",
        ".tsx": "typescript-react",
        ".java": "java",
        ".c": "c",
        ".cpp": "cpp",
        ".h": "c-header",
        ".hpp": "cpp-header",
        ".go": "go",
        ".rs": "rust",
        ".rb": "ruby",
        ".php": "php",
        ".swift": "swift",
        ".kt": "kotlin",
        ".scala": "scala",
        ".r": "r",
        ".sql": "sql",
        ".sh": "bash",
        ".bash": "bash",
        ".zsh": "zsh",
        ".ps1": "powershell",
        ".yaml": "yaml",
        ".yml": "yaml",
        ".json": "json",
        ".xml": "xml",
        ".html": "html",
        ".css": "css",
        ".scss": "scss",
        ".less": "less",
    }

    @classmethod
    def read(cls, file_path: Path) -> str:
        language = cls.LANGUAGE_MAP.get(file_path.suffix.lower(), "unknown")
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        return f"[Language: {language}]\n[File: {file_path.name}]\n\n{content}"


class WebReader(DocumentReader):
    """Reader for web content via URL."""

    @staticmethod
    def read(url: str) -> str:
        try:
            import trafilatura
        except ImportError:
            import subprocess

            subprocess.check_call(
                [sys.executable, "-m", "pip", "install", "trafilatura", "--break-system-packages", "-q"]
            )
            import trafilatura

        downloaded = trafilatura.fetch_url(url)
        if not downloaded:
            raise RuntimeError(f"Failed to fetch URL: {url}")

        content = trafilatura.extract(downloaded, include_tables=True, include_links=True)
        if not content:
            raise RuntimeError(f"Failed to extract content from URL: {url}")

        # Try to get metadata
        metadata = trafilatura.extract_metadata(downloaded)
        title = getattr(metadata, "title", None) if metadata else None
        author = getattr(metadata, "author", None) if metadata else None
        date = getattr(metadata, "date", None) if metadata else None

        header_parts = [f"[Source: {url}]"]
        if title:
            header_parts.append(f"[Title: {title}]")
        if author:
            header_parts.append(f"[Author: {author}]")
        if date:
            header_parts.append(f"[Date: {date}]")

        return "\n".join(header_parts) + "\n\n" + content


class YouTubeTranscriptReader(DocumentReader):
    """Reader for YouTube video transcripts."""

    # Patterns that identify YouTube URLs
    _PATTERNS = [
        re.compile(r"(?:youtube\.com/watch\?.*v=|youtu\.be/)([\w-]{11})"),
    ]

    @classmethod
    def extract_video_id(cls, url: str) -> str | None:
        for pat in cls._PATTERNS:
            m = pat.search(url)
            if m:
                return m.group(1)
        return None

    @classmethod
    def can_handle(cls, url: str) -> bool:
        return cls.extract_video_id(url) is not None

    @staticmethod
    def read(url: str) -> str:  # type: ignore[override]
        video_id = YouTubeTranscriptReader.extract_video_id(url)
        if not video_id:
            raise RuntimeError(f"Could not extract YouTube video ID from: {url}")

        try:
            from youtube_transcript_api import YouTubeTranscriptApi
        except ImportError:
            raise RuntimeError(
                "YouTube transcript reader requires youtube-transcript-api: "
                "pip install youtube-transcript-api"
            )

        try:
            transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
            # Prefer manually created transcripts, fall back to auto-generated
            try:
                transcript = transcript_list.find_manually_created_transcript(["en"])
            except Exception:
                transcript = transcript_list.find_generated_transcript(["en"])

            entries = transcript.fetch()
            # Build readable transcript
            parts = [f"[Source: {url}]", f"[Format: YouTube Transcript]", f"[Video ID: {video_id}]", ""]
            for entry in entries:
                text = entry.get("text", entry) if isinstance(entry, dict) else str(entry)
                parts.append(str(text))

            return "\n".join(parts)
        except Exception as e:
            raise RuntimeError(f"Failed to fetch YouTube transcript: {e}")


class DataReader(DocumentReader):
    """Reader for structured data files (JSON, YAML, CSV, XML)."""

    DATA_EXTENSIONS = {".json", ".yaml", ".yml", ".csv", ".xml"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        suffix = file_path.suffix.lower()
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()

        format_name = {
            ".json": "JSON",
            ".yaml": "YAML",
            ".yml": "YAML",
            ".csv": "CSV",
            ".xml": "XML",
        }.get(suffix, "Data")

        structure_desc = cls._describe_structure(content, suffix)
        return f"[Format: {format_name}]\n[File: {file_path.name}]\n\n{structure_desc}\n\n--- Raw Content ---\n{content}"

    @classmethod
    def _describe_structure(cls, content: str, suffix: str) -> str:
        """Generate a structural description of the data."""
        try:
            if suffix == ".json":
                import json

                data = json.loads(content)
                return cls._describe_json_structure(data)
            elif suffix in {".yaml", ".yml"}:
                try:
                    import yaml

                    data = yaml.safe_load(content)
                    return cls._describe_json_structure(data)
                except ImportError:
                    return "[YAML parsing unavailable]"
            elif suffix == ".csv":
                return cls._describe_csv_structure(content)
            elif suffix == ".xml":
                return cls._describe_xml_structure(content)
        except Exception as e:
            return f"[Structure analysis failed: {e}]"
        return ""

    @classmethod
    def _describe_json_structure(cls, data: Any, depth: int = 0, max_depth: int = 3) -> str:
        """Describe the structure of JSON/YAML data."""
        indent = "  " * depth
        if depth >= max_depth:
            return f"{indent}..."

        if isinstance(data, dict):
            if not data:
                return f"{indent}{{}}"
            lines = [f"{indent}Object with {len(data)} keys:"]
            for key, value in list(data.items())[:10]:
                value_type = type(value).__name__
                if isinstance(value, dict):
                    lines.append(f"{indent}  - {key}: object ({len(value)} keys)")
                elif isinstance(value, list):
                    lines.append(f"{indent}  - {key}: array ({len(value)} items)")
                else:
                    lines.append(f"{indent}  - {key}: {value_type}")
            if len(data) > 10:
                lines.append(f"{indent}  ... and {len(data) - 10} more keys")
            return "\n".join(lines)
        elif isinstance(data, list):
            if not data:
                return f"{indent}[]"
            lines = [f"{indent}Array with {len(data)} items"]
            if data:
                first = data[0]
                if isinstance(first, dict):
                    lines.append(f"{indent}  Item type: object with keys: {list(first.keys())[:5]}")
                else:
                    lines.append(f"{indent}  Item type: {type(first).__name__}")
            return "\n".join(lines)
        else:
            return f"{indent}{type(data).__name__}: {str(data)[:100]}"

    @classmethod
    def _describe_csv_structure(cls, content: str) -> str:
        """Describe CSV structure."""
        lines = content.strip().split("\n")
        if not lines:
            return "[Empty CSV]"

        header_line = lines[0]
        columns = [col.strip().strip('"').strip("'") for col in header_line.split(",")]

        desc = [f"CSV with {len(columns)} columns and {len(lines) - 1} data rows"]
        desc.append(f"Columns: {', '.join(columns[:10])}")
        if len(columns) > 10:
            desc.append(f"  ... and {len(columns) - 10} more columns")

        return "\n".join(desc)

    @classmethod
    def _describe_xml_structure(cls, content: str) -> str:
        """Describe XML structure (basic)."""
        import re

        root_match = re.search(r"<(\w+)[>\s]", content)
        root_tag = root_match.group(1) if root_match else "unknown"

        tag_counts: dict[str, int] = {}
        for tag in re.findall(r"<(\w+)[>\s]", content):
            tag_counts[tag] = tag_counts.get(tag, 0) + 1

        top_tags = sorted(tag_counts.items(), key=lambda x: -x[1])[:10]

        desc = [f"XML document with root element: <{root_tag}>"]
        desc.append(f"Top elements: {', '.join(f'{t}({c})' for t, c in top_tags)}")

        return "\n".join(desc)


class PDFReader(DocumentReader):
    @staticmethod
    def read(file_path: Path) -> str:
        try:
            import pdfplumber
        except ImportError:
            import subprocess

            subprocess.check_call(
                [sys.executable, "-m", "pip", "install", "pdfplumber", "--break-system-packages", "-q"]
            )
            import pdfplumber

        text_parts: list[str] = []
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"[Page {i + 1}]\n{page_text}")
        return "\n\n".join(text_parts)


class ImageReader(DocumentReader):
    """Reader for images using OCR."""

    IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        try:
            from PIL import Image
        except ImportError:
            import subprocess

            subprocess.check_call(
                [sys.executable, "-m", "pip", "install", "Pillow", "--break-system-packages", "-q"]
            )
            from PIL import Image

        try:
            import pytesseract
        except ImportError:
            import subprocess

            subprocess.check_call(
                [sys.executable, "-m", "pip", "install", "pytesseract", "--break-system-packages", "-q"]
            )
            import pytesseract

        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            if not text.strip():
                return f"[Image: {file_path.name}]\n[No text detected via OCR]"
            return f"[Image: {file_path.name}]\n[OCR Extracted Text]\n\n{text}"
        except Exception as e:
            return f"[Image: {file_path.name}]\n[OCR failed: {e}]"


class AudioReader(DocumentReader):
    """Reader for audio files using speech-to-text."""

    AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac", ".wma"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        try:
            import whisper
        except ImportError:
            import subprocess

            subprocess.check_call(
                [sys.executable, "-m", "pip", "install", "openai-whisper", "--break-system-packages", "-q"]
            )
            import whisper

        try:
            model = whisper.load_model("base")
            result = model.transcribe(str(file_path))
            text = result.get("text", "")
            if not text.strip():
                return f"[Audio: {file_path.name}]\n[No speech detected]"
            return f"[Audio: {file_path.name}]\n[Transcription]\n\n{text}"
        except Exception as e:
            return f"[Audio: {file_path.name}]\n[Transcription failed: {e}]"


class VideoReader(DocumentReader):
    """Reader for video files - extracts audio and transcribes."""

    VIDEO_EXTENSIONS = {".mp4", ".avi", ".mov", ".mkv", ".webm", ".wmv", ".flv"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        try:
            from moviepy.editor import VideoFileClip
        except ImportError:
            import subprocess

            subprocess.check_call(
                [sys.executable, "-m", "pip", "install", "moviepy", "--break-system-packages", "-q"]
            )
            from moviepy.editor import VideoFileClip

        try:
            import whisper
        except ImportError:
            import subprocess

            subprocess.check_call(
                [sys.executable, "-m", "pip", "install", "openai-whisper", "--break-system-packages", "-q"]
            )
            import whisper

        import tempfile

        try:
            # Extract audio from video
            video = VideoFileClip(str(file_path))
            duration = video.duration

            with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp_audio:
                audio_path = tmp_audio.name
                video.audio.write_audiofile(audio_path, verbose=False, logger=None)
                video.close()

            # Transcribe the audio
            model = whisper.load_model("base")
            result = model.transcribe(audio_path)
            text = result.get("text", "")

            # Clean up temp file
            import os
            os.unlink(audio_path)

            if not text.strip():
                return f"[Video: {file_path.name}]\n[Duration: {duration:.1f}s]\n[No speech detected]"

            return f"[Video: {file_path.name}]\n[Duration: {duration:.1f}s]\n[Transcription]\n\n{text}"

        except Exception as e:
            return f"[Video: {file_path.name}]\n[Transcription failed: {e}]"


class DocxReader(DocumentReader):
    """Reader for Microsoft Word .docx files."""

    EXTENSIONS = {".docx"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        try:
            import docx
        except ImportError:
            raise RuntimeError(
                "DOCX reader requires python-docx: pip install hexis[readers]"
            )

        doc = docx.Document(file_path)
        parts: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                parts.append("\t".join(cells))
        return f"[Format: DOCX]\n[File: {file_path.name}]\n\n" + "\n".join(parts)


class RtfReader(DocumentReader):
    """Reader for RTF files."""

    EXTENSIONS = {".rtf"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        try:
            from striprtf.striprtf import rtf_to_text
        except ImportError:
            raise RuntimeError(
                "RTF reader requires striprtf: pip install hexis[readers]"
            )

        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            raw = f.read()
        text = rtf_to_text(raw)
        return f"[Format: RTF]\n[File: {file_path.name}]\n\n{text}"


class LatexReader(DocumentReader):
    """Reader for LaTeX .tex and .bib files (regex-based, no deps)."""

    EXTENSIONS = {".tex", ".bib"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()

        suffix = file_path.suffix.lower()
        if suffix == ".bib":
            return cls._read_bib(content, file_path)
        return cls._read_tex(content, file_path)

    @classmethod
    def _read_tex(cls, content: str, file_path: Path) -> str:
        # Strip non-content environments
        for env in ("figure", "tikzpicture", "table", "lstlisting"):
            content = re.sub(
                rf"\\begin\{{{env}\}}.*?\\end\{{{env}\}}",
                "",
                content,
                flags=re.DOTALL,
            )
        # Strip comments
        content = re.sub(r"(?m)%.*$", "", content)
        # Strip \command{...} but keep content inside braces for text commands
        content = re.sub(r"\\(?:textbf|textit|emph|underline)\{([^}]*)\}", r"\1", content)
        # Strip remaining commands (keep braced content)
        content = re.sub(r"\\[a-zA-Z]+\*?(?:\[[^\]]*\])*\{([^}]*)\}", r"\1", content)
        # Strip standalone commands
        content = re.sub(r"\\[a-zA-Z]+\*?(?:\[[^\]]*\])*", "", content)
        # Clean up braces
        content = content.replace("{", "").replace("}", "")
        # Collapse whitespace
        content = re.sub(r"\n{3,}", "\n\n", content).strip()
        return f"[Format: LaTeX]\n[File: {file_path.name}]\n\n{content}"

    @classmethod
    def _read_bib(cls, content: str, file_path: Path) -> str:
        entries: list[str] = []
        for match in re.finditer(
            r"@\w+\{([^,]+),\s*(.*?)\n\}", content, flags=re.DOTALL
        ):
            key = match.group(1).strip()
            body = match.group(2)
            fields: dict[str, str] = {}
            for fm in re.finditer(r"(\w+)\s*=\s*\{([^}]*)\}", body):
                fields[fm.group(1).lower()] = fm.group(2).strip()
            parts = [f"[{key}]"]
            for field_name in ("author", "title", "year", "abstract"):
                if field_name in fields:
                    parts.append(f"  {field_name.title()}: {fields[field_name]}")
            entries.append("\n".join(parts))
        return (
            f"[Format: BibTeX]\n[File: {file_path.name}]\n[Entries: {len(entries)}]\n\n"
            + "\n\n".join(entries)
        )


class EmailReader(DocumentReader):
    """Reader for .eml and .mbox email files (stdlib only)."""

    EXTENSIONS = {".eml", ".mbox"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        suffix = file_path.suffix.lower()
        if suffix == ".mbox":
            return cls._read_mbox(file_path)
        return cls._read_eml(file_path)

    @classmethod
    def _read_eml(cls, file_path: Path) -> str:
        import email
        import email.policy

        with open(file_path, "rb") as f:
            msg = email.message_from_bytes(f.read(), policy=email.policy.default)
        return cls._format_message(msg)

    @classmethod
    def _read_mbox(cls, file_path: Path) -> str:
        import mailbox

        mbox = mailbox.mbox(str(file_path))
        parts: list[str] = []
        for i, msg in enumerate(mbox):
            parts.append(f"--- Message {i + 1} ---")
            parts.append(cls._format_message(msg))
        if not parts:
            return f"[Email]\n[File: {file_path.name}]\n[Empty mailbox]"
        return "\n\n".join(parts)

    @classmethod
    def _format_message(cls, msg: Any) -> str:
        subject = str(msg.get("Subject", "(no subject)"))
        from_addr = str(msg.get("From", ""))
        to_addr = str(msg.get("To", ""))
        date = str(msg.get("Date", ""))

        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                ct = part.get_content_type()
                if ct == "text/plain":
                    payload = part.get_payload(decode=True)
                    if payload:
                        body = payload.decode("utf-8", errors="replace")
                        break
            if not body:
                for part in msg.walk():
                    ct = part.get_content_type()
                    if ct == "text/html":
                        payload = part.get_payload(decode=True)
                        if payload:
                            body = re.sub(r"<[^>]+>", "", payload.decode("utf-8", errors="replace"))
                            break
        else:
            payload = msg.get_payload(decode=True)
            if payload:
                body = payload.decode("utf-8", errors="replace")
            elif isinstance(msg.get_payload(), str):
                body = msg.get_payload()

        header = f"[Email]\n[Subject: {subject}]\n[From: {from_addr}]\n[To: {to_addr}]\n[Date: {date}]"
        return f"{header}\n\n{body.strip()}"


class EpubReader(DocumentReader):
    """Reader for EPUB e-book files."""

    EXTENSIONS = {".epub"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        try:
            import ebooklib
            from ebooklib import epub
        except ImportError:
            raise RuntimeError(
                "EPUB reader requires ebooklib and beautifulsoup4: pip install hexis[readers]"
            )

        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise RuntimeError(
                "EPUB reader requires beautifulsoup4: pip install hexis[readers]"
            )

        book = epub.read_epub(str(file_path), options={"ignore_ncx": True})
        title = book.get_metadata("DC", "title")
        title_str = title[0][0] if title else file_path.stem
        author = book.get_metadata("DC", "creator")
        author_str = author[0][0] if author else "Unknown"

        parts: list[str] = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_body_content(), "html.parser")
            text = soup.get_text(separator="\n", strip=True)
            if text.strip():
                parts.append(text)

        header = f"[Format: EPUB]\n[Title: {title_str}]\n[Author: {author_str}]"
        return f"{header}\n\n" + "\n\n".join(parts)


class PptxReader(DocumentReader):
    """Reader for PowerPoint .pptx files."""

    EXTENSIONS = {".pptx"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            raise RuntimeError(
                "PPTX reader requires python-pptx: pip install hexis[readers]"
            )

        prs = Presentation(str(file_path))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                parts.append(f"[Slide {i + 1}]\n" + "\n".join(slide_texts))

        header = f"[Format: PPTX]\n[File: {file_path.name}]\n[Slides: {len(prs.slides)}]"
        return f"{header}\n\n" + "\n\n".join(parts)


class XlsxReader(DocumentReader):
    """Reader for Excel .xlsx/.xls files."""

    EXTENSIONS = {".xlsx", ".xls"}

    MAX_ROWS_PER_SHEET = 500

    @classmethod
    def read(cls, file_path: Path) -> str:
        try:
            import openpyxl
        except ImportError:
            raise RuntimeError(
                "XLSX reader requires openpyxl: pip install hexis[readers]"
            )

        wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
        parts: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[str] = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= cls.MAX_ROWS_PER_SHEET:
                    rows.append(f"... (truncated at {cls.MAX_ROWS_PER_SHEET} rows)")
                    break
                cells = [str(c) if c is not None else "" for c in row]
                rows.append("\t".join(cells))
            parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()

        header = f"[Format: XLSX]\n[File: {file_path.name}]\n[Sheets: {len(wb.sheetnames)}]"
        return f"{header}\n\n" + "\n\n".join(parts)


class NotebookReader(DocumentReader):
    """Reader for Jupyter .ipynb notebooks (stdlib json only)."""

    EXTENSIONS = {".ipynb"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            nb = json.loads(f.read())

        cells = nb.get("cells", [])
        kernel = (
            nb.get("metadata", {}).get("kernelspec", {}).get("language", "python")
        )

        parts: list[str] = []
        for cell in cells:
            cell_type = cell.get("cell_type", "")
            source = "".join(cell.get("source", []))
            if not source.strip():
                continue
            if cell_type == "markdown":
                parts.append(source)
            elif cell_type == "code":
                parts.append(f"```{kernel}\n{source}\n```")

        header = f"[Format: Jupyter Notebook]\n[File: {file_path.name}]\n[Cells: {len(cells)}]"
        return f"{header}\n\n" + "\n\n".join(parts)


class RssReader(DocumentReader):
    """Reader for RSS/Atom feeds (URL-based, not file-based)."""

    @classmethod
    def read(cls, url: str) -> str:  # type: ignore[override]
        try:
            import feedparser
        except ImportError:
            raise RuntimeError(
                "RSS reader requires feedparser: pip install hexis[readers]"
            )

        feed = feedparser.parse(url)
        if not feed.entries:
            return ""

        feed_title = feed.feed.get("title", url)
        parts: list[str] = [f"[Format: RSS/Atom]\n[Feed: {feed_title}]\n[Entries: {len(feed.entries)}]"]
        for entry in feed.entries:
            title = entry.get("title", "(no title)")
            link = entry.get("link", "")
            summary = entry.get("summary", "")
            published = entry.get("published", "")
            # Strip HTML from summary
            summary = re.sub(r"<[^>]+>", "", summary)
            entry_parts = [f"## {title}"]
            if published:
                entry_parts.append(f"Published: {published}")
            if link:
                entry_parts.append(f"Link: {link}")
            if summary:
                entry_parts.append(summary)
            parts.append("\n".join(entry_parts))

        return "\n\n".join(parts)


def get_reader(file_path: Path) -> DocumentReader:
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return PDFReader()
    if suffix in [".md", ".markdown"]:
        return MarkdownReader()
    if suffix in DocxReader.EXTENSIONS:
        return DocxReader()
    if suffix in RtfReader.EXTENSIONS:
        return RtfReader()
    if suffix in LatexReader.EXTENSIONS:
        return LatexReader()
    if suffix in EmailReader.EXTENSIONS:
        return EmailReader()
    if suffix in EpubReader.EXTENSIONS:
        return EpubReader()
    if suffix in PptxReader.EXTENSIONS:
        return PptxReader()
    if suffix in XlsxReader.EXTENSIONS:
        return XlsxReader()
    if suffix in NotebookReader.EXTENSIONS:
        return NotebookReader()
    if suffix in DataReader.DATA_EXTENSIONS:
        return DataReader()
    if suffix in ImageReader.IMAGE_EXTENSIONS:
        return ImageReader()
    if suffix in AudioReader.AUDIO_EXTENSIONS:
        return AudioReader()
    if suffix in VideoReader.VIDEO_EXTENSIONS:
        return VideoReader()
    if suffix in CodeReader.LANGUAGE_MAP:
        return CodeReader()
    return TextReader()

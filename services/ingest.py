#!/usr/bin/env python3
"""
Hexis Universal Ingestion Pipeline

Implements the ingestion flow described in ToDo/ingest.md.
"""

from __future__ import annotations

import argparse
import asyncio
import hashlib
import json
import os
import re
import sys
import time
from dataclasses import dataclass, field
from datetime import datetime, timezone
from enum import Enum
from pathlib import Path
from typing import Any, Callable, Optional
from uuid import UUID

from core.cognitive_memory_api import (
    CognitiveMemorySync,
    MemoryInput as ApiMemoryInput,
    MemoryType as ApiMemoryType,
    RelationshipType,
)

# =========================================================================
# CONFIGURATION
# =========================================================================


class IngestionMode(str, Enum):
    """Ingestion mode taxonomy:

    FAST   -- chunk, appraise, extract facts (default)
    SLOW   -- RLM conscious reading per chunk
    HYBRID -- fast triage then selective RLM on high-signal chunks
    """

    FAST = "fast"
    SLOW = "slow"
    HYBRID = "hybrid"


@dataclass
class Config:
    """Pipeline configuration."""

    # LLM Settings -- a fully-resolved llm_config dict (from resolve_llm_config / load_llm_config)
    llm_config: dict[str, Any] | None = None

    # Database Settings (unified -- pass a DSN string)
    dsn: str | None = None
    # Legacy DB fields (used when dsn is None -- standalone CLI)
    db_host: str = "localhost"
    db_port: int = 43815
    db_name: str = "hexis_memory"
    db_user: str = "postgres"
    db_password: str = "password"

    # Mode
    mode: IngestionMode = IngestionMode.FAST

    # Internal threshold: docs <= this word count get per-section appraisal
    deep_max_words: int = 2000

    # Chunking
    max_section_chars: int = 2000
    chunk_overlap: int = 200

    # Extraction
    max_facts_per_section: int = 20
    min_confidence_threshold: float = 0.6
    skip_sections: list[str] = field(
        default_factory=lambda: ["references", "bibliography", "acknowledgments", "appendix"]
    )

    # Persistence overrides
    min_importance_floor: float | None = None
    permanent: bool = False

    # Source trust override
    base_trust: float | None = None

    # Processing
    verbose: bool = True
    log: Optional[Callable[[str], None]] = None
    cancel_check: Optional[Callable[[], bool]] = None


# =========================================================================
# DATA STRUCTURES
# =========================================================================


@dataclass
class Section:
    title: str
    content: str
    index: int


@dataclass
class DocumentInfo:
    title: str
    source_type: str
    content_hash: str
    word_count: int
    path: str
    file_type: str


@dataclass
class Appraisal:
    valence: float = 0.0
    arousal: float = 0.3
    primary_emotion: str = "neutral"
    intensity: float = 0.0
    goal_relevance: list[dict[str, Any]] = field(default_factory=list)
    worldview_tension: float = 0.0
    curiosity: float = 0.0
    summary: str = ""

    def to_state_payload(self, source: str = "ingest") -> dict[str, Any]:
        return {
            "valence": self.valence,
            "arousal": self.arousal,
            "primary_emotion": self.primary_emotion,
            "intensity": self.intensity,
            "source": source,
        }


@dataclass
class Extraction:
    content: str
    category: str
    confidence: float
    importance: float
    why: str | None = None
    connections: list[str] = field(default_factory=list)
    supports: str | None = None
    contradicts: str | None = None
    concepts: list[str] = field(default_factory=list)


@dataclass
class IngestionMetrics:
    """Metrics collected during ingestion for observability."""

    source_type: str = ""
    source_size_bytes: int = 0
    word_count: int = 0
    mode: str = ""
    appraisal_valence: float = 0.0
    appraisal_arousal: float = 0.0
    appraisal_emotion: str = "neutral"
    appraisal_intensity: float = 0.0
    extraction_count: int = 0
    dedup_count: int = 0
    memory_count: int = 0
    llm_calls: int = 0
    duration_seconds: float = 0.0
    errors: list[str] = field(default_factory=list)
    start_time: float = field(default_factory=lambda: 0.0)


# =========================================================================
# HELPERS
# =========================================================================


def _emit(config: Config, message: str) -> None:
    if config.log:
        config.log(message)
    else:
        print(message)


def _should_cancel(config: Config) -> bool:
    if config.cancel_check:
        try:
            return bool(config.cancel_check())
        except Exception:
            return False
    return False


def _hash_text(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


def _word_count(text: str) -> int:
    return len(re.findall(r"\b\w+\b", text))


def _normalize_mode(mode: IngestionMode | str | None) -> IngestionMode:
    if isinstance(mode, IngestionMode):
        return mode
    raw = str(mode or "fast").strip().lower()
    # Legacy modes all collapse to FAST
    if raw in ("auto", "standard", "deep", "shallow", "archive"):
        return IngestionMode.FAST
    for item in IngestionMode:
        if raw == item.value:
            return item
    return IngestionMode.FAST


def _decay_rate_for_intensity(intensity: float, base: float = 0.01) -> float:
    if intensity < 0.1:
        return base * 3.0
    if intensity < 0.3:
        return base * 1.5
    if intensity > 0.6:
        return base * 0.5
    return base


def _infer_source_type(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    if suffix in {".pdf", ".md", ".markdown", ".txt", ".text", ".rtf", ".docx", ".tex", ".bib", ".epub"}:
        return "document"
    if suffix == ".pptx":
        return "presentation"
    if suffix in {".xlsx", ".xls"}:
        return "spreadsheet"
    if suffix in {".eml", ".mbox"}:
        return "email"
    if suffix == ".ipynb":
        return "code"
    if suffix in CodeReader.LANGUAGE_MAP:
        return "code"
    if suffix in {".json", ".yaml", ".yml", ".csv", ".xml"}:
        return "data"
    if suffix in ImageReader.IMAGE_EXTENSIONS:
        return "image"
    if suffix in AudioReader.AUDIO_EXTENSIONS:
        return "audio"
    if suffix in VideoReader.VIDEO_EXTENSIONS:
        return "video"
    return "document"


def _extract_title(content: str, file_path: Path) -> str:
    # Try markdown header
    header_match = re.search(r"^#\s+(.+)$", content, re.MULTILINE)
    if header_match:
        return header_match.group(1).strip()
    # Try first non-empty line
    for line in content.splitlines():
        if line.strip():
            return line.strip()[:120]
    return file_path.stem


# =========================================================================
# DOCUMENT READERS
# =========================================================================


class DocumentReader:
    @staticmethod
    def read(file_path: Path) -> str:
        raise NotImplementedError


class MarkdownReader(DocumentReader):
    @staticmethod
    def read(file_path: Path) -> str:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()


class TextReader(DocumentReader):
    @staticmethod
    def read(file_path: Path) -> str:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()


class CodeReader(DocumentReader):
    LANGUAGE_MAP = {
        ".py": "python",
        ".js": "javascript",
        ".ts": "typescript",
        ".jsx": "javascript-react",
        ".tsx": "typescript-react",
        ".java": "java",
        ".c": "c",
        ".cpp": "cpp",
        ".h": "c-header",
        ".hpp": "cpp-header",
        ".go": "go",
        ".rs": "rust",
        ".rb": "ruby",
        ".php": "php",
        ".swift": "swift",
        ".kt": "kotlin",
        ".scala": "scala",
        ".r": "r",
        ".sql": "sql",
        ".sh": "bash",
        ".bash": "bash",
        ".zsh": "zsh",
        ".ps1": "powershell",
        ".yaml": "yaml",
        ".yml": "yaml",
        ".json": "json",
        ".xml": "xml",
        ".html": "html",
        ".css": "css",
        ".scss": "scss",
        ".less": "less",
    }

    @classmethod
    def read(cls, file_path: Path) -> str:
        language = cls.LANGUAGE_MAP.get(file_path.suffix.lower(), "unknown")
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        return f"[Language: {language}]\n[File: {file_path.name}]\n\n{content}"


class WebReader(DocumentReader):
    """Reader for web content via URL."""

    @staticmethod
    def read(url: str) -> str:
        try:
            import trafilatura
        except ImportError:
            import subprocess

            subprocess.check_call(
                [sys.executable, "-m", "pip", "install", "trafilatura", "--break-system-packages", "-q"]
            )
            import trafilatura

        downloaded = trafilatura.fetch_url(url)
        if not downloaded:
            raise RuntimeError(f"Failed to fetch URL: {url}")

        content = trafilatura.extract(downloaded, include_tables=True, include_links=True)
        if not content:
            raise RuntimeError(f"Failed to extract content from URL: {url}")

        # Try to get metadata
        metadata = trafilatura.extract_metadata(downloaded)
        title = getattr(metadata, "title", None) if metadata else None
        author = getattr(metadata, "author", None) if metadata else None
        date = getattr(metadata, "date", None) if metadata else None

        header_parts = [f"[Source: {url}]"]
        if title:
            header_parts.append(f"[Title: {title}]")
        if author:
            header_parts.append(f"[Author: {author}]")
        if date:
            header_parts.append(f"[Date: {date}]")

        return "\n".join(header_parts) + "\n\n" + content


class YouTubeTranscriptReader(DocumentReader):
    """Reader for YouTube video transcripts."""

    # Patterns that identify YouTube URLs
    _PATTERNS = [
        re.compile(r"(?:youtube\.com/watch\?.*v=|youtu\.be/)([\w-]{11})"),
    ]

    @classmethod
    def extract_video_id(cls, url: str) -> str | None:
        for pat in cls._PATTERNS:
            m = pat.search(url)
            if m:
                return m.group(1)
        return None

    @classmethod
    def can_handle(cls, url: str) -> bool:
        return cls.extract_video_id(url) is not None

    @staticmethod
    def read(url: str) -> str:  # type: ignore[override]
        video_id = YouTubeTranscriptReader.extract_video_id(url)
        if not video_id:
            raise RuntimeError(f"Could not extract YouTube video ID from: {url}")

        try:
            from youtube_transcript_api import YouTubeTranscriptApi
        except ImportError:
            raise RuntimeError(
                "YouTube transcript reader requires youtube-transcript-api: "
                "pip install youtube-transcript-api"
            )

        try:
            transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
            # Prefer manually created transcripts, fall back to auto-generated
            try:
                transcript = transcript_list.find_manually_created_transcript(["en"])
            except Exception:
                transcript = transcript_list.find_generated_transcript(["en"])

            entries = transcript.fetch()
            # Build readable transcript
            parts = [f"[Source: {url}]", f"[Format: YouTube Transcript]", f"[Video ID: {video_id}]", ""]
            for entry in entries:
                text = entry.get("text", entry) if isinstance(entry, dict) else str(entry)
                parts.append(str(text))

            return "\n".join(parts)
        except Exception as e:
            raise RuntimeError(f"Failed to fetch YouTube transcript: {e}")


class DataReader(DocumentReader):
    """Reader for structured data files (JSON, YAML, CSV, XML)."""

    DATA_EXTENSIONS = {".json", ".yaml", ".yml", ".csv", ".xml"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        suffix = file_path.suffix.lower()
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()

        format_name = {
            ".json": "JSON",
            ".yaml": "YAML",
            ".yml": "YAML",
            ".csv": "CSV",
            ".xml": "XML",
        }.get(suffix, "Data")

        structure_desc = cls._describe_structure(content, suffix)
        return f"[Format: {format_name}]\n[File: {file_path.name}]\n\n{structure_desc}\n\n--- Raw Content ---\n{content}"

    @classmethod
    def _describe_structure(cls, content: str, suffix: str) -> str:
        """Generate a structural description of the data."""
        try:
            if suffix == ".json":
                import json

                data = json.loads(content)
                return cls._describe_json_structure(data)
            elif suffix in {".yaml", ".yml"}:
                try:
                    import yaml

                    data = yaml.safe_load(content)
                    return cls._describe_json_structure(data)
                except ImportError:
                    return "[YAML parsing unavailable]"
            elif suffix == ".csv":
                return cls._describe_csv_structure(content)
            elif suffix == ".xml":
                return cls._describe_xml_structure(content)
        except Exception as e:
            return f"[Structure analysis failed: {e}]"
        return ""

    @classmethod
    def _describe_json_structure(cls, data: Any, depth: int = 0, max_depth: int = 3) -> str:
        """Describe the structure of JSON/YAML data."""
        indent = "  " * depth
        if depth >= max_depth:
            return f"{indent}..."

        if isinstance(data, dict):
            if not data:
                return f"{indent}{{}}"
            lines = [f"{indent}Object with {len(data)} keys:"]
            for key, value in list(data.items())[:10]:
                value_type = type(value).__name__
                if isinstance(value, dict):
                    lines.append(f"{indent}  - {key}: object ({len(value)} keys)")
                elif isinstance(value, list):
                    lines.append(f"{indent}  - {key}: array ({len(value)} items)")
                else:
                    lines.append(f"{indent}  - {key}: {value_type}")
            if len(data) > 10:
                lines.append(f"{indent}  ... and {len(data) - 10} more keys")
            return "\n".join(lines)
        elif isinstance(data, list):
            if not data:
                return f"{indent}[]"
            lines = [f"{indent}Array with {len(data)} items"]
            if data:
                first = data[0]
                if isinstance(first, dict):
                    lines.append(f"{indent}  Item type: object with keys: {list(first.keys())[:5]}")
                else:
                    lines.append(f"{indent}  Item type: {type(first).__name__}")
            return "\n".join(lines)
        else:
            return f"{indent}{type(data).__name__}: {str(data)[:100]}"

    @classmethod
    def _describe_csv_structure(cls, content: str) -> str:
        """Describe CSV structure."""
        lines = content.strip().split("\n")
        if not lines:
            return "[Empty CSV]"

        header_line = lines[0]
        columns = [col.strip().strip('"').strip("'") for col in header_line.split(",")]

        desc = [f"CSV with {len(columns)} columns and {len(lines) - 1} data rows"]
        desc.append(f"Columns: {', '.join(columns[:10])}")
        if len(columns) > 10:
            desc.append(f"  ... and {len(columns) - 10} more columns")

        return "\n".join(desc)

    @classmethod
    def _describe_xml_structure(cls, content: str) -> str:
        """Describe XML structure (basic)."""
        import re

        root_match = re.search(r"<(\w+)[>\s]", content)
        root_tag = root_match.group(1) if root_match else "unknown"

        tag_counts: dict[str, int] = {}
        for tag in re.findall(r"<(\w+)[>\s]", content):
            tag_counts[tag] = tag_counts.get(tag, 0) + 1

        top_tags = sorted(tag_counts.items(), key=lambda x: -x[1])[:10]

        desc = [f"XML document with root element: <{root_tag}>"]
        desc.append(f"Top elements: {', '.join(f'{t}({c})' for t, c in top_tags)}")

        return "\n".join(desc)


class PDFReader(DocumentReader):
    @staticmethod
    def read(file_path: Path) -> str:
        try:
            import pdfplumber
        except ImportError:
            import subprocess

            subprocess.check_call(
                [sys.executable, "-m", "pip", "install", "pdfplumber", "--break-system-packages", "-q"]
            )
            import pdfplumber

        text_parts: list[str] = []
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"[Page {i + 1}]\n{page_text}")
        return "\n\n".join(text_parts)


class ImageReader(DocumentReader):
    """Reader for images using OCR."""

    IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        try:
            from PIL import Image
        except ImportError:
            import subprocess

            subprocess.check_call(
                [sys.executable, "-m", "pip", "install", "Pillow", "--break-system-packages", "-q"]
            )
            from PIL import Image

        try:
            import pytesseract
        except ImportError:
            import subprocess

            subprocess.check_call(
                [sys.executable, "-m", "pip", "install", "pytesseract", "--break-system-packages", "-q"]
            )
            import pytesseract

        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            if not text.strip():
                return f"[Image: {file_path.name}]\n[No text detected via OCR]"
            return f"[Image: {file_path.name}]\n[OCR Extracted Text]\n\n{text}"
        except Exception as e:
            return f"[Image: {file_path.name}]\n[OCR failed: {e}]"


class AudioReader(DocumentReader):
    """Reader for audio files using speech-to-text."""

    AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac", ".wma"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        try:
            import whisper
        except ImportError:
            import subprocess

            subprocess.check_call(
                [sys.executable, "-m", "pip", "install", "openai-whisper", "--break-system-packages", "-q"]
            )
            import whisper

        try:
            model = whisper.load_model("base")
            result = model.transcribe(str(file_path))
            text = result.get("text", "")
            if not text.strip():
                return f"[Audio: {file_path.name}]\n[No speech detected]"
            return f"[Audio: {file_path.name}]\n[Transcription]\n\n{text}"
        except Exception as e:
            return f"[Audio: {file_path.name}]\n[Transcription failed: {e}]"


class VideoReader(DocumentReader):
    """Reader for video files - extracts audio and transcribes."""

    VIDEO_EXTENSIONS = {".mp4", ".avi", ".mov", ".mkv", ".webm", ".wmv", ".flv"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        try:
            from moviepy.editor import VideoFileClip
        except ImportError:
            import subprocess

            subprocess.check_call(
                [sys.executable, "-m", "pip", "install", "moviepy", "--break-system-packages", "-q"]
            )
            from moviepy.editor import VideoFileClip

        try:
            import whisper
        except ImportError:
            import subprocess

            subprocess.check_call(
                [sys.executable, "-m", "pip", "install", "openai-whisper", "--break-system-packages", "-q"]
            )
            import whisper

        import tempfile

        try:
            # Extract audio from video
            video = VideoFileClip(str(file_path))
            duration = video.duration

            with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp_audio:
                audio_path = tmp_audio.name
                video.audio.write_audiofile(audio_path, verbose=False, logger=None)
                video.close()

            # Transcribe the audio
            model = whisper.load_model("base")
            result = model.transcribe(audio_path)
            text = result.get("text", "")

            # Clean up temp file
            import os
            os.unlink(audio_path)

            if not text.strip():
                return f"[Video: {file_path.name}]\n[Duration: {duration:.1f}s]\n[No speech detected]"

            return f"[Video: {file_path.name}]\n[Duration: {duration:.1f}s]\n[Transcription]\n\n{text}"

        except Exception as e:
            return f"[Video: {file_path.name}]\n[Transcription failed: {e}]"


class DocxReader(DocumentReader):
    """Reader for Microsoft Word .docx files."""

    EXTENSIONS = {".docx"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        try:
            import docx
        except ImportError:
            raise RuntimeError(
                "DOCX reader requires python-docx: pip install hexis[readers]"
            )

        doc = docx.Document(file_path)
        parts: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                parts.append("\t".join(cells))
        return f"[Format: DOCX]\n[File: {file_path.name}]\n\n" + "\n".join(parts)


class RtfReader(DocumentReader):
    """Reader for RTF files."""

    EXTENSIONS = {".rtf"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        try:
            from striprtf.striprtf import rtf_to_text
        except ImportError:
            raise RuntimeError(
                "RTF reader requires striprtf: pip install hexis[readers]"
            )

        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            raw = f.read()
        text = rtf_to_text(raw)
        return f"[Format: RTF]\n[File: {file_path.name}]\n\n{text}"


class LatexReader(DocumentReader):
    """Reader for LaTeX .tex and .bib files (regex-based, no deps)."""

    EXTENSIONS = {".tex", ".bib"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()

        suffix = file_path.suffix.lower()
        if suffix == ".bib":
            return cls._read_bib(content, file_path)
        return cls._read_tex(content, file_path)

    @classmethod
    def _read_tex(cls, content: str, file_path: Path) -> str:
        # Strip non-content environments
        for env in ("figure", "tikzpicture", "table", "lstlisting"):
            content = re.sub(
                rf"\\begin\{{{env}\}}.*?\\end\{{{env}\}}",
                "",
                content,
                flags=re.DOTALL,
            )
        # Strip comments
        content = re.sub(r"(?m)%.*$", "", content)
        # Strip \command{...} but keep content inside braces for text commands
        content = re.sub(r"\\(?:textbf|textit|emph|underline)\{([^}]*)\}", r"\1", content)
        # Strip remaining commands (keep braced content)
        content = re.sub(r"\\[a-zA-Z]+\*?(?:\[[^\]]*\])*\{([^}]*)\}", r"\1", content)
        # Strip standalone commands
        content = re.sub(r"\\[a-zA-Z]+\*?(?:\[[^\]]*\])*", "", content)
        # Clean up braces
        content = content.replace("{", "").replace("}", "")
        # Collapse whitespace
        content = re.sub(r"\n{3,}", "\n\n", content).strip()
        return f"[Format: LaTeX]\n[File: {file_path.name}]\n\n{content}"

    @classmethod
    def _read_bib(cls, content: str, file_path: Path) -> str:
        entries: list[str] = []
        for match in re.finditer(
            r"@\w+\{([^,]+),\s*(.*?)\n\}", content, flags=re.DOTALL
        ):
            key = match.group(1).strip()
            body = match.group(2)
            fields: dict[str, str] = {}
            for fm in re.finditer(r"(\w+)\s*=\s*\{([^}]*)\}", body):
                fields[fm.group(1).lower()] = fm.group(2).strip()
            parts = [f"[{key}]"]
            for field_name in ("author", "title", "year", "abstract"):
                if field_name in fields:
                    parts.append(f"  {field_name.title()}: {fields[field_name]}")
            entries.append("\n".join(parts))
        return (
            f"[Format: BibTeX]\n[File: {file_path.name}]\n[Entries: {len(entries)}]\n\n"
            + "\n\n".join(entries)
        )


class EmailReader(DocumentReader):
    """Reader for .eml and .mbox email files (stdlib only)."""

    EXTENSIONS = {".eml", ".mbox"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        suffix = file_path.suffix.lower()
        if suffix == ".mbox":
            return cls._read_mbox(file_path)
        return cls._read_eml(file_path)

    @classmethod
    def _read_eml(cls, file_path: Path) -> str:
        import email
        import email.policy

        with open(file_path, "rb") as f:
            msg = email.message_from_bytes(f.read(), policy=email.policy.default)
        return cls._format_message(msg)

    @classmethod
    def _read_mbox(cls, file_path: Path) -> str:
        import mailbox

        mbox = mailbox.mbox(str(file_path))
        parts: list[str] = []
        for i, msg in enumerate(mbox):
            parts.append(f"--- Message {i + 1} ---")
            parts.append(cls._format_message(msg))
        if not parts:
            return f"[Email]\n[File: {file_path.name}]\n[Empty mailbox]"
        return "\n\n".join(parts)

    @classmethod
    def _format_message(cls, msg: Any) -> str:
        subject = str(msg.get("Subject", "(no subject)"))
        from_addr = str(msg.get("From", ""))
        to_addr = str(msg.get("To", ""))
        date = str(msg.get("Date", ""))

        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                ct = part.get_content_type()
                if ct == "text/plain":
                    payload = part.get_payload(decode=True)
                    if payload:
                        body = payload.decode("utf-8", errors="replace")
                        break
            if not body:
                for part in msg.walk():
                    ct = part.get_content_type()
                    if ct == "text/html":
                        payload = part.get_payload(decode=True)
                        if payload:
                            body = re.sub(r"<[^>]+>", "", payload.decode("utf-8", errors="replace"))
                            break
        else:
            payload = msg.get_payload(decode=True)
            if payload:
                body = payload.decode("utf-8", errors="replace")
            elif isinstance(msg.get_payload(), str):
                body = msg.get_payload()

        header = f"[Email]\n[Subject: {subject}]\n[From: {from_addr}]\n[To: {to_addr}]\n[Date: {date}]"
        return f"{header}\n\n{body.strip()}"


class EpubReader(DocumentReader):
    """Reader for EPUB e-book files."""

    EXTENSIONS = {".epub"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        try:
            import ebooklib
            from ebooklib import epub
        except ImportError:
            raise RuntimeError(
                "EPUB reader requires ebooklib and beautifulsoup4: pip install hexis[readers]"
            )

        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise RuntimeError(
                "EPUB reader requires beautifulsoup4: pip install hexis[readers]"
            )

        book = epub.read_epub(str(file_path), options={"ignore_ncx": True})
        title = book.get_metadata("DC", "title")
        title_str = title[0][0] if title else file_path.stem
        author = book.get_metadata("DC", "creator")
        author_str = author[0][0] if author else "Unknown"

        parts: list[str] = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_body_content(), "html.parser")
            text = soup.get_text(separator="\n", strip=True)
            if text.strip():
                parts.append(text)

        header = f"[Format: EPUB]\n[Title: {title_str}]\n[Author: {author_str}]"
        return f"{header}\n\n" + "\n\n".join(parts)


class PptxReader(DocumentReader):
    """Reader for PowerPoint .pptx files."""

    EXTENSIONS = {".pptx"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            raise RuntimeError(
                "PPTX reader requires python-pptx: pip install hexis[readers]"
            )

        prs = Presentation(str(file_path))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                parts.append(f"[Slide {i + 1}]\n" + "\n".join(slide_texts))

        header = f"[Format: PPTX]\n[File: {file_path.name}]\n[Slides: {len(prs.slides)}]"
        return f"{header}\n\n" + "\n\n".join(parts)


class XlsxReader(DocumentReader):
    """Reader for Excel .xlsx/.xls files."""

    EXTENSIONS = {".xlsx", ".xls"}

    MAX_ROWS_PER_SHEET = 500

    @classmethod
    def read(cls, file_path: Path) -> str:
        try:
            import openpyxl
        except ImportError:
            raise RuntimeError(
                "XLSX reader requires openpyxl: pip install hexis[readers]"
            )

        wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
        parts: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[str] = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= cls.MAX_ROWS_PER_SHEET:
                    rows.append(f"... (truncated at {cls.MAX_ROWS_PER_SHEET} rows)")
                    break
                cells = [str(c) if c is not None else "" for c in row]
                rows.append("\t".join(cells))
            parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()

        header = f"[Format: XLSX]\n[File: {file_path.name}]\n[Sheets: {len(wb.sheetnames)}]"
        return f"{header}\n\n" + "\n\n".join(parts)


class NotebookReader(DocumentReader):
    """Reader for Jupyter .ipynb notebooks (stdlib json only)."""

    EXTENSIONS = {".ipynb"}

    @classmethod
    def read(cls, file_path: Path) -> str:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            nb = json.loads(f.read())

        cells = nb.get("cells", [])
        kernel = (
            nb.get("metadata", {}).get("kernelspec", {}).get("language", "python")
        )

        parts: list[str] = []
        for cell in cells:
            cell_type = cell.get("cell_type", "")
            source = "".join(cell.get("source", []))
            if not source.strip():
                continue
            if cell_type == "markdown":
                parts.append(source)
            elif cell_type == "code":
                parts.append(f"```{kernel}\n{source}\n```")

        header = f"[Format: Jupyter Notebook]\n[File: {file_path.name}]\n[Cells: {len(cells)}]"
        return f"{header}\n\n" + "\n\n".join(parts)


class RssReader(DocumentReader):
    """Reader for RSS/Atom feeds (URL-based, not file-based)."""

    @classmethod
    def read(cls, url: str) -> str:  # type: ignore[override]
        try:
            import feedparser
        except ImportError:
            raise RuntimeError(
                "RSS reader requires feedparser: pip install hexis[readers]"
            )

        feed = feedparser.parse(url)
        if not feed.entries:
            return ""

        feed_title = feed.feed.get("title", url)
        parts: list[str] = [f"[Format: RSS/Atom]\n[Feed: {feed_title}]\n[Entries: {len(feed.entries)}]"]
        for entry in feed.entries:
            title = entry.get("title", "(no title)")
            link = entry.get("link", "")
            summary = entry.get("summary", "")
            published = entry.get("published", "")
            # Strip HTML from summary
            summary = re.sub(r"<[^>]+>", "", summary)
            entry_parts = [f"## {title}"]
            if published:
                entry_parts.append(f"Published: {published}")
            if link:
                entry_parts.append(f"Link: {link}")
            if summary:
                entry_parts.append(summary)
            parts.append("\n".join(entry_parts))

        return "\n\n".join(parts)


def get_reader(file_path: Path) -> DocumentReader:
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return PDFReader()
    if suffix in [".md", ".markdown"]:
        return MarkdownReader()
    if suffix in DocxReader.EXTENSIONS:
        return DocxReader()
    if suffix in RtfReader.EXTENSIONS:
        return RtfReader()
    if suffix in LatexReader.EXTENSIONS:
        return LatexReader()
    if suffix in EmailReader.EXTENSIONS:
        return EmailReader()
    if suffix in EpubReader.EXTENSIONS:
        return EpubReader()
    if suffix in PptxReader.EXTENSIONS:
        return PptxReader()
    if suffix in XlsxReader.EXTENSIONS:
        return XlsxReader()
    if suffix in NotebookReader.EXTENSIONS:
        return NotebookReader()
    if suffix in DataReader.DATA_EXTENSIONS:
        return DataReader()
    if suffix in ImageReader.IMAGE_EXTENSIONS:
        return ImageReader()
    if suffix in AudioReader.AUDIO_EXTENSIONS:
        return AudioReader()
    if suffix in VideoReader.VIDEO_EXTENSIONS:
        return VideoReader()
    if suffix in CodeReader.LANGUAGE_MAP:
        return CodeReader()
    return TextReader()


# =========================================================================
# SECTIONING
# =========================================================================


class Sectioner:
    def __init__(self, max_chars: int = 2000, overlap: int = 200):
        self.max_chars = max_chars
        self.overlap = overlap

    def split(self, content: str, file_path: Path) -> list[Section]:
        suffix = file_path.suffix.lower()
        if suffix in [".md", ".markdown"]:
            return self._split_markdown(content)
        if suffix == ".pptx":
            return self._split_on_delimiter(content, r"\[Slide \d+\]")
        if suffix in {".xlsx", ".xls"}:
            return self._split_on_delimiter(content, r"\[Sheet: [^\]]+\]")
        if suffix == ".ipynb":
            return self._split_notebook(content)
        if suffix in {".eml", ".mbox"}:
            return self._split_on_delimiter(content, r"--- Message \d+ ---")
        return self._split_text(content)

    def _split_markdown(self, content: str) -> list[Section]:
        header_pattern = r"^(#{1,6}\s+.+)$"
        parts = re.split(header_pattern, content, flags=re.MULTILINE)
        sections: list[Section] = []
        current_title = "Introduction"
        current_content = ""
        for part in parts:
            if re.match(header_pattern, part or ""):
                if current_content.strip():
                    sections.append(Section(title=current_title, content=current_content.strip(), index=len(sections)))
                current_title = part.strip().lstrip("# ").strip() or current_title
                current_content = ""
            else:
                current_content += part
        if current_content.strip():
            sections.append(Section(title=current_title, content=current_content.strip(), index=len(sections)))
        if not sections:
            return [Section(title="Document", content=content, index=0)]
        return sections

    def _split_text(self, content: str) -> list[Section]:
        if len(content) <= self.max_chars:
            return [Section(title="Section 1", content=content, index=0)]
        chunks: list[str] = []
        paragraphs = content.split("\n\n")
        current = ""
        for para in paragraphs:
            if len(current) + len(para) + 2 <= self.max_chars:
                current += para + "\n\n"
                continue
            if current:
                chunks.append(current.strip())
            if len(para) > self.max_chars:
                sentences = re.split(r"(?<=[.!?])\s+", para)
                current = ""
                for sentence in sentences:
                    if len(current) + len(sentence) <= self.max_chars:
                        current += sentence + " "
                    else:
                        if current:
                            chunks.append(current.strip())
                        current = sentence + " "
            else:
                current = para + "\n\n"
        if current.strip():
            chunks.append(current.strip())
        if self.overlap > 0 and len(chunks) > 1:
            overlapped: list[str] = []
            for i, chunk in enumerate(chunks):
                if i > 0:
                    prev_overlap = chunks[i - 1][-self.overlap :]
                    chunk = f"...{prev_overlap}\n\n{chunk}"
                overlapped.append(chunk)
            chunks = overlapped
        return [Section(title=f"Section {i + 1}", content=chunk, index=i) for i, chunk in enumerate(chunks)]

    def _split_on_delimiter(self, content: str, pattern: str) -> list[Section]:
        """Split content on a regex delimiter pattern, keeping the delimiter with its section."""
        parts = re.split(f"({pattern})", content)
        sections: list[Section] = []
        current_title = "Header"
        current_content = ""
        for part in parts:
            if re.match(pattern, part):
                if current_content.strip():
                    sections.append(Section(title=current_title, content=current_content.strip(), index=len(sections)))
                current_title = part.strip().strip("[]").strip("-").strip()
                current_content = ""
            else:
                current_content += part
        if current_content.strip():
            sections.append(Section(title=current_title, content=current_content.strip(), index=len(sections)))
        if not sections:
            return [Section(title="Document", content=content, index=0)]
        return sections

    def _split_notebook(self, content: str) -> list[Section]:
        """Split notebook content on cell boundaries (triple-backtick blocks and text)."""
        parts = re.split(r"(```\w*\n.*?```)", content, flags=re.DOTALL)
        sections: list[Section] = []
        for part in parts:
            text = part.strip()
            if not text:
                continue
            if text.startswith("```"):
                title = f"Code Cell {len(sections) + 1}"
            else:
                title = f"Cell {len(sections) + 1}"
            sections.append(Section(title=title, content=text, index=len(sections)))
        if not sections:
            return [Section(title="Notebook", content=content, index=0)]
        return sections


# =========================================================================
# LLM CLIENT
# =========================================================================


class IngestLLM:
    """Thin LLM wrapper for the ingestion pipeline.

    Uses core.llm.chat_completion() under the hood, supporting all
    configured providers (OpenAI, Anthropic, Codex, Gemini, etc.).
    """

    def __init__(self, config: Config):
        from core.llm import normalize_llm_config

        self._cfg = normalize_llm_config(config.llm_config)
        self.call_count = 0

    def complete(self, messages: list[dict[str, str]], temperature: float = 0.3) -> str:
        """Sync completion -- uses asyncio.run() for the standalone CLI path."""
        return asyncio.run(self.acomplete(messages, temperature=temperature))

    async def acomplete(self, messages: list[dict[str, str]], temperature: float = 0.3) -> str:
        """Async completion via core.llm."""
        from core.llm import chat_completion
        from core.usage import record_llm_usage

        self.call_count += 1
        result = await chat_completion(
            provider=self._cfg["provider"],
            model=self._cfg["model"],
            endpoint=self._cfg.get("endpoint"),
            api_key=self._cfg.get("api_key"),
            messages=messages,
            temperature=temperature,
            max_tokens=1200,
            auth_mode=self._cfg.get("auth_mode"),
        )
        asyncio.ensure_future(record_llm_usage(
            provider=self._cfg["provider"],
            model=self._cfg["model"],
            raw_response=result.get("raw"),
            source="ingest",
        ))
        return result.get("content", "")

    def complete_json(self, messages: list[dict[str, str]], temperature: float = 0.2) -> dict[str, Any]:
        text = self.complete(messages, temperature=temperature)
        return self._parse_json(text)

    async def acomplete_json(self, messages: list[dict[str, str]], temperature: float = 0.2) -> dict[str, Any]:
        """Async version of complete_json() for use with asyncio.gather()."""
        text = await self.acomplete(messages, temperature=temperature)
        return self._parse_json(text)

    @staticmethod
    def _parse_json(text: str) -> dict[str, Any]:
        json_text = text.strip()
        match = re.search(r"```(?:json)?\s*(.*?)\s*```", json_text, re.DOTALL)
        if match:
            json_text = match.group(1).strip()
        if not json_text.startswith("{"):
            start = json_text.find("{")
            if start != -1:
                json_text = json_text[start:]
        try:
            data = json.loads(json_text)
        except Exception:
            return {}
        return data if isinstance(data, dict) else {}


# =========================================================================
# APPRAISAL + EXTRACTION
# =========================================================================


class Appraiser:
    def __init__(self, llm: IngestLLM):
        self.llm = llm

    @staticmethod
    def _build_messages(content: str, context: dict[str, Any]) -> list[dict[str, str]]:
        system = (
            "You are Hexis' subconscious appraisal system."
            " Provide a brief, honest emotional assessment of the content."
            " If you feel nothing, say so and keep intensity low."
            " Return STRICT JSON only."
        )
        user = (
            "CONTENT SAMPLE:\n"
            f"{content}\n\n"
            "CONTEXT (JSON):\n"
            f"{json.dumps(context)[:8000]}\n\n"
            "Return JSON with keys:"
            " valence (-1..1), arousal (0..1), primary_emotion (string), intensity (0..1),"
            " goal_relevance (array of {goal, strength}), worldview_tension (0..1), curiosity (0..1),"
            " summary (2-3 sentences)."
        )
        return [{"role": "system", "content": system}, {"role": "user", "content": user}]

    @staticmethod
    def _parse(raw: dict[str, Any]) -> Appraisal:
        return Appraisal(
            valence=float(raw.get("valence", 0.0) or 0.0),
            arousal=float(raw.get("arousal", 0.3) or 0.3),
            primary_emotion=str(raw.get("primary_emotion", "neutral") or "neutral"),
            intensity=float(raw.get("intensity", 0.0) or 0.0),
            goal_relevance=list(raw.get("goal_relevance", []) or []),
            worldview_tension=float(raw.get("worldview_tension", 0.0) or 0.0),
            curiosity=float(raw.get("curiosity", 0.0) or 0.0),
            summary=str(raw.get("summary", "") or ""),
        )

    def appraise(self, *, content: str, context: dict[str, Any], mode: IngestionMode) -> Appraisal:
        msgs = self._build_messages(content, context)
        raw = self.llm.complete_json(msgs, temperature=0.2)
        return self._parse(raw)

    async def aappraise(self, *, content: str, context: dict[str, Any], mode: IngestionMode) -> Appraisal:
        """Async version for use with asyncio.gather()."""
        msgs = self._build_messages(content, context)
        raw = await self.llm.acomplete_json(msgs, temperature=0.2)
        return self._parse(raw)


class KnowledgeExtractor:
    def __init__(self, llm: IngestLLM):
        self.llm = llm

    @staticmethod
    def _build_messages(
        section: Section, doc: DocumentInfo, appraisal: Appraisal, mode: IngestionMode, max_items: int,
    ) -> list[dict[str, str]]:
        system = (
            "You extract standalone knowledge worth remembering."
            " Be selective. Return STRICT JSON only."
        )
        if doc.source_type == "code":
            guidance = (
                "Focus on what the code does, key interfaces, behaviors, patterns,"
                " and any important constraints or dependencies."
            )
        elif doc.source_type == "data":
            guidance = (
                "Describe the schema, key fields, relationships, and notable values or patterns."
            )
        else:
            guidance = (
                "Extract facts, claims, definitions, procedures, insights, and statistics."
            )
        user = (
            f"DOCUMENT: {doc.title}\n"
            f"SECTION: {section.title}\n"
            f"MODE: {mode.value}\n\n"
            "APPRAISAL:\n"
            f"{json.dumps(appraisal.__dict__, ensure_ascii=False)}\n\n"
            "CONTENT:\n"
            f"{section.content}\n\n"
            f"{guidance}\n\n"
            "Return JSON with key 'items' as an array of objects:\n"
            "  {content, category, confidence, importance, why, connections, supports, contradicts, concepts}\n"
            "  - concepts: array of key concept/entity names this knowledge is an instance of\n"
            "Keep at most "
            + str(max_items)
            + " items."
        )
        return [{"role": "system", "content": system}, {"role": "user", "content": user}]

    @staticmethod
    def _parse(raw: dict[str, Any], max_items: int) -> list[Extraction]:
        items = raw.get("items") if isinstance(raw, dict) else None
        if not isinstance(items, list):
            return []
        out: list[Extraction] = []
        for item in items[:max_items]:
            if not isinstance(item, dict):
                continue
            content = str(item.get("content", "") or "").strip()
            if not content:
                continue
            out.append(
                Extraction(
                    content=content,
                    category=str(item.get("category", "fact") or "fact"),
                    confidence=float(item.get("confidence", 0.5) or 0.5),
                    importance=float(item.get("importance", 0.5) or 0.5),
                    why=str(item.get("why", "") or "") or None,
                    connections=[str(c).strip() for c in (item.get("connections") or []) if str(c).strip()],
                    supports=item.get("supports"),
                    contradicts=item.get("contradicts"),
                    concepts=[str(c).strip() for c in (item.get("concepts") or []) if str(c).strip()],
                )
            )
        return out

    def extract(
        self,
        *,
        section: Section,
        doc: DocumentInfo,
        appraisal: Appraisal,
        mode: IngestionMode,
        max_items: int,
    ) -> list[Extraction]:
        msgs = self._build_messages(section, doc, appraisal, mode, max_items)
        raw = self.llm.complete_json(msgs, temperature=0.3)
        return self._parse(raw, max_items)

    async def aextract(
        self,
        *,
        section: Section,
        doc: DocumentInfo,
        appraisal: Appraisal,
        mode: IngestionMode,
        max_items: int,
    ) -> list[Extraction]:
        """Async version for use with asyncio.gather()."""
        msgs = self._build_messages(section, doc, appraisal, mode, max_items)
        raw = await self.llm.acomplete_json(msgs, temperature=0.3)
        return self._parse(raw, max_items)


# =========================================================================
# STORAGE
# =========================================================================


class MemoryStore:
    def __init__(self, config: Config):
        self.config = config
        self.client: CognitiveMemorySync | None = None

    def connect(self) -> None:
        if self.client is not None:
            return
        if self.config.dsn:
            dsn = self.config.dsn
        else:
            dsn = (
                f"postgresql://{self.config.db_user}:{self.config.db_password}"
                f"@{self.config.db_host}:{self.config.db_port}/{self.config.db_name}"
            )
        self.client = CognitiveMemorySync.connect(dsn, min_size=1, max_size=5)

    def close(self) -> None:
        if self.client is not None:
            self.client.close()
            self.client = None

    def _exec(self, sql: str, *params: Any) -> Any:
        assert self.client is not None
        async def _run():
            async with self.client._async._pool.acquire() as conn:
                return await conn.execute(sql, *params)

        return self.client._loop.run_until_complete(_run())

    def _fetchval(self, sql: str, *params: Any) -> Any:
        assert self.client is not None
        async def _run():
            async with self.client._async._pool.acquire() as conn:
                return await conn.fetchval(sql, *params)

        return self.client._loop.run_until_complete(_run())

    def has_receipt(self, content_hash: str) -> bool:
        if self.client is None:
            self.connect()
        assert self.client is not None
        try:
            receipts = self.client.get_ingestion_receipts(content_hash, [content_hash])
        except Exception:
            return False
        return bool(receipts)

    def set_affective_state(self, appraisal: Appraisal) -> None:
        if self.client is None:
            self.connect()
        payload = json.dumps(appraisal.to_state_payload(source="ingest"))
        try:
            self._fetchval("SELECT set_current_affective_state($1::jsonb)", payload)
        except Exception:
            pass

    def create_encounter_memory(
        self,
        *,
        text: str,
        source: dict[str, Any],
        emotional_valence: float,
        context: dict[str, Any] | None,
        importance: float,
    ) -> str:
        if self.client is None:
            self.connect()
        assert self.client is not None
        memory_id = self.client.remember(
            text,
            type=ApiMemoryType.EPISODIC,
            importance=importance,
            emotional_valence=emotional_valence,
            context=context,
            source_attribution=source,
        )
        return str(memory_id)

    def create_semantic_memory(
        self,
        *,
        content: str,
        confidence: float,
        category: str,
        related_concepts: list[str],
        source: dict[str, Any],
        importance: float,
        trust: float | None,
    ) -> str:
        if self.client is None:
            self.connect()
        payload_sources = json.dumps([source])
        return str(
            self._fetchval(
                "SELECT create_semantic_memory($1::text,$2::float,$3::text[],$4::text[],$5::jsonb,$6::float,$7::jsonb,$8::float)",
                content,
                confidence,
                [category],
                related_concepts,
                payload_sources,
                importance,
                json.dumps(source),
                trust,
            )
        )

    def add_source(self, memory_id: str, source: dict[str, Any]) -> None:
        if self.client is None:
            self.connect()
        assert self.client is not None
        self.client._loop.run_until_complete(
            self.client._async.add_source(UUID(memory_id), source)
        )

    def add_evidence(
        self,
        memory_id: str,
        stance: str,
        source: dict[str, Any],
        note: str | None = None,
        evidence_memory_id: str | None = None,
        context: str = "ingest",
    ) -> dict[str, Any]:
        """Attach evidence to an existing memory through the DB-owned belief
        revision policy (db/59): source merge, SUPPORTS/CONTRADICTS edge,
        calibrated confidence update, and an audit row. Returns the revision
        result ({prior, posterior, applied, reason, ...})."""
        if self.client is None:
            self.connect()
        raw = self._fetchval(
            "SELECT add_memory_evidence($1::uuid, $2::text, $3::jsonb, $4::text, $5::uuid, $6::text)",
            memory_id,
            stance,
            json.dumps(source),
            note,
            evidence_memory_id,
            context,
        )
        parsed = json.loads(raw) if isinstance(raw, str) else raw
        return parsed if isinstance(parsed, dict) else {}

    def link_concept(self, memory_id: str, concept: str, strength: float = 1.0) -> None:
        """Link a memory to a concept in the knowledge graph."""
        if self.client is None:
            self.connect()
        self._fetchval(
            "SELECT link_memory_to_concept($1::uuid, $2::text, $3::float)",
            memory_id,
            concept,
            strength,
        )

    def link_concepts_batch(self, pairs: list[tuple[str, str]], strength: float = 1.0) -> None:
        """Link multiple (memory_id, concept) pairs in a single batch call."""
        if not pairs:
            return
        if self.client is None:
            self.connect()
        assert self.client is not None

        async def _run():
            async with self.client._async._pool.acquire() as conn:
                await conn.executemany(
                    "SELECT link_memory_to_concept($1::uuid, $2::text, $3::float)",
                    [(mid, concept, strength) for mid, concept in pairs],
                )

        self.client._loop.run_until_complete(_run())

    def connect_memories_batch(self, edges: list[tuple[str, str, "RelationshipType", float]]) -> None:
        """Create multiple memory relationships in a single batch call."""
        if not edges:
            return
        if self.client is None:
            self.connect()
        assert self.client is not None
        from core.cognitive_memory_api import RelationshipInput
        rels = [
            RelationshipInput(
                from_id=UUID(from_id),
                to_id=UUID(to_id),
                relationship_type=rel_type,
                confidence=conf,
            )
            for from_id, to_id, rel_type, conf in edges
        ]
        self.client.connect_batch(rels)

    def prefetch_embeddings(self, texts: list[str]) -> int:
        """Pre-warm embedding cache for a batch of texts.

        Calls the SQL ``prefetch_embeddings()`` function which batches HTTP
        requests to the embedding service (default batch size 8) and caches
        results.  Subsequent ``recall_similar_semantic`` / ``create_semantic_memory``
        calls for the same content become cache hits.
        """
        if not texts:
            return 0
        if self.client is None:
            self.connect()
        return self._fetchval("SELECT prefetch_embeddings($1::text[])", texts) or 0

    def recall_similar_semantic(self, query: str, limit: int = 5):
        if self.client is None:
            self.connect()
        assert self.client is not None
        return self.client.recall(
            query,
            limit=limit,
            memory_types=[ApiMemoryType.SEMANTIC],
        ).memories

    def recall_similar(self, query: str, memory_types: list[str], limit: int = 5):
        """Recall nearest memories of the given types (list result)."""
        if self.client is None:
            self.connect()
        assert self.client is not None
        return self.client.recall(
            query,
            limit=limit,
            memory_types=[ApiMemoryType(t) for t in memory_types],
        ).memories

    def route_extractions(self, extractions: list, min_confidence: float) -> list:
        """Route extractions through the DB dedup/related/create policy
        (db/41 ingest_route_extractions): config-driven thresholds + one batched
        nearest-neighbor search. Returns a per-extraction plan with 'index',
        'decision' (duplicate|related|create) and 'matched_memory_id'."""
        if self.client is None:
            self.connect()
        payload = json.dumps([
            {"content": ext.content, "confidence": ext.confidence}
            for ext in extractions
        ])
        raw = self._fetchval(
            "SELECT ingest_route_extractions($1::jsonb, $2::float)", payload, min_confidence
        )
        plan = json.loads(raw) if isinstance(raw, str) else raw
        return plan or []

    def route_texts(self, items: list[tuple[str, float]], min_confidence: float = 0.0) -> list:
        """Route bare (content, confidence) pairs through the same DB
        dedup/related/create policy as route_extractions."""
        if not items:
            return []
        if self.client is None:
            self.connect()
        payload = json.dumps([
            {"content": content, "confidence": confidence}
            for content, confidence in items
        ])
        raw = self._fetchval(
            "SELECT ingest_route_extractions($1::jsonb, $2::float)", payload, min_confidence
        )
        plan = json.loads(raw) if isinstance(raw, str) else raw
        return plan or []

    def connect_memories(self, from_id: str, to_id: str, relationship: RelationshipType, confidence: float = 0.8) -> None:
        if self.client is None:
            self.connect()
        assert self.client is not None
        self.client.connect_memories(
            from_id,
            to_id,
            relationship,
            confidence=confidence,
        )

    def update_decay_rate(self, memory_id: str, decay_rate: float) -> None:
        if self.client is None:
            self.connect()
        try:
            self._exec("UPDATE memories SET decay_rate = $1 WHERE id = $2::uuid", decay_rate, memory_id)
        except Exception:
            pass

    def fetch_appraisal_context(self) -> dict[str, Any]:
        if self.client is None:
            self.connect()
        try:
            raw = self._fetchval(
                """
                SELECT jsonb_build_object(
                    'emotional_state', get_current_affective_state(),
                    'goals', get_goals_snapshot(),
                    'worldview', get_worldview_context(),
                    'recent_memories', get_recent_context(5)
                )
                """
            )
            if isinstance(raw, str):
                return json.loads(raw)
            if isinstance(raw, dict):
                return raw
        except Exception:
            return {}
        return {}

    def store_metrics(self, metrics: "IngestionMetrics") -> None:
        """Store ingestion metrics for observability."""
        if self.client is None:
            self.connect()
        try:
            self._exec(
                """
                INSERT INTO ingestion_metrics (
                    source_type, source_size_bytes, word_count, mode,
                    appraisal_valence, appraisal_arousal, appraisal_emotion, appraisal_intensity,
                    extraction_count, dedup_count, memory_count, llm_calls,
                    duration_seconds, errors
                ) VALUES (
                    $1, $2, $3, $4, $5, $6, $7, $8, $9, $10, $11, $12, $13, $14::jsonb
                )
                """,
                metrics.source_type,
                metrics.source_size_bytes,
                metrics.word_count,
                metrics.mode,
                metrics.appraisal_valence,
                metrics.appraisal_arousal,
                metrics.appraisal_emotion,
                metrics.appraisal_intensity,
                metrics.extraction_count,
                metrics.dedup_count,
                metrics.memory_count,
                metrics.llm_calls,
                metrics.duration_seconds,
                json.dumps(metrics.errors),
            )
        except Exception:
            pass  # Don't fail ingestion due to metrics storage

    def check_archived_for_query(self, query: str, threshold: float = 0.75, limit: int = 5) -> list[dict[str, Any]]:
        """Check if archived content matches a query."""
        if self.client is None:
            self.connect()
        try:
            rows = self._fetchval(
                """
                SELECT jsonb_agg(jsonb_build_object(
                    'memory_id', memory_id,
                    'content_hash', content_hash,
                    'title', title,
                    'similarity', similarity,
                    'source_path', source_path
                ))
                FROM check_archived_for_query($1, $2, $3)
                """,
                query,
                threshold,
                limit,
            )
            if not rows:
                return []
            result = json.loads(rows) if isinstance(rows, str) else rows
            return result if result else []
        except Exception:
            return []

    def mark_archived_processed(self, memory_id: str) -> bool:
        """Mark an archived memory as processed."""
        if self.client is None:
            self.connect()
        try:
            result = self._fetchval(
                "SELECT mark_archived_as_processed($1::uuid)",
                memory_id,
            )
            return bool(result)
        except Exception:
            return False


# =========================================================================
# INGESTION PIPELINE
# =========================================================================


class IngestionPipeline:
    SUPPORTED_EXTENSIONS = (
        # Documents
        {".md", ".markdown", ".txt", ".text", ".pdf"}
        | DocxReader.EXTENSIONS
        | RtfReader.EXTENSIONS
        | LatexReader.EXTENSIONS
        | EmailReader.EXTENSIONS
        | EpubReader.EXTENSIONS
        | PptxReader.EXTENSIONS
        | XlsxReader.EXTENSIONS
        | NotebookReader.EXTENSIONS
        # Code
        | set(CodeReader.LANGUAGE_MAP.keys())
        # Data
        | DataReader.DATA_EXTENSIONS
        # Media
        | ImageReader.IMAGE_EXTENSIONS
        | AudioReader.AUDIO_EXTENSIONS
        | VideoReader.VIDEO_EXTENSIONS
    )

    def __init__(self, config: Config):
        self.config = config
        self.config.mode = _normalize_mode(self.config.mode)
        self.sectioner = Sectioner(config.max_section_chars, config.chunk_overlap)
        self.llm = IngestLLM(config)
        self.appraiser = Appraiser(self.llm)
        self.extractor = KnowledgeExtractor(self.llm)
        self.store = MemoryStore(config)
        self.stats = {"files_processed": 0, "memories_created": 0, "errors": 0}

    def ingest_file(self, file_path: Path) -> int:
        # Initialize metrics tracking
        metrics = IngestionMetrics(start_time=time.time())

        if _should_cancel(self.config):
            raise RuntimeError("Ingestion cancelled")
        if not file_path.exists():
            _emit(self.config, f"File not found: {file_path}")
            return 0
        if file_path.suffix.lower() not in self.SUPPORTED_EXTENSIONS:
            _emit(self.config, f"Unsupported file type: {file_path.suffix}")
            return 0

        if self.config.verbose:
            _emit(self.config, f"\nProcessing: {file_path}")

        # Track LLM calls at start
        llm_calls_start = self.llm.call_count

        reader = get_reader(file_path)
        try:
            content = reader.read(file_path)
            metrics.source_size_bytes = len(content.encode("utf-8"))
        except Exception as exc:
            _emit(self.config, f"  Error reading file: {exc}")
            self.stats["errors"] += 1
            metrics.errors.append(str(exc))
            return 0

        title = _extract_title(content, file_path)
        words = _word_count(content)
        mode = self.config.mode
        source_type = _infer_source_type(file_path)
        content_hash = _hash_text(content)

        # Update metrics
        metrics.word_count = words
        metrics.mode = mode.value
        metrics.source_type = source_type

        doc = DocumentInfo(
            title=title,
            source_type=source_type,
            content_hash=content_hash,
            word_count=words,
            path=str(file_path),
            file_type=file_path.suffix.lower(),
        )

        if self.store.has_receipt(content_hash):
            if self.config.verbose:
                _emit(self.config, f"  Already ingested (hash={content_hash[:8]}...). Skipping.")
            return 0

        sections = self.sectioner.split(content, file_path)

        if self.config.verbose:
            _emit(self.config, f"  Mode: {mode.value} | Words: {words} | Sections: {len(sections)}")

        # Slow/hybrid mode: delegate to RLM-based ingestion
        if mode in (IngestionMode.SLOW, IngestionMode.HYBRID):
            count = self._run_rlm_ingest(mode, doc, sections, metrics, llm_calls_start)
            return count

        # FAST mode: small docs (<=deep_max_words) get per-section appraisal;
        # larger docs get a single doc-level appraisal. All sections processed.
        base_context = self._build_appraisal_context(doc)
        use_deep = words <= self.config.deep_max_words

        overall_appraisal = None
        if not use_deep:
            sample = self._sample_content(content)
            overall_appraisal = self.appraiser.appraise(content=sample, context=base_context, mode=mode)
            self.store.set_affective_state(overall_appraisal)
            metrics.appraisal_valence = overall_appraisal.valence
            metrics.appraisal_arousal = overall_appraisal.arousal
            metrics.appraisal_emotion = overall_appraisal.primary_emotion
            metrics.appraisal_intensity = overall_appraisal.intensity

        encounter_id = self._create_encounter_memory(doc, overall_appraisal, mode)

        created_ids: list[str] = []
        total_extractions = 0
        dedup_count = 0

        # -- Phase 1: parallel LLM extraction --
        active_sections = [s for s in sections if not self._skip_section(s.title)]

        if _should_cancel(self.config):
            raise RuntimeError("Ingestion cancelled")

        max_items = self.config.max_facts_per_section

        if use_deep:
            # Small docs: each section gets its own appraisal + extraction
            async def _parallel_deep() -> list[tuple[Appraisal, list[Extraction]]]:
                async def _appraise_and_extract(s: Section) -> tuple[Appraisal, list[Extraction]]:
                    sample = self._sample_content(s.content)
                    apr = await self.appraiser.aappraise(content=sample, context=base_context, mode=mode)
                    exts = await self.extractor.aextract(
                        section=s, doc=doc, appraisal=apr, mode=mode, max_items=max_items,
                    )
                    return apr, exts
                return await asyncio.gather(*[_appraise_and_extract(s) for s in active_sections])

            deep_results = asyncio.run(_parallel_deep())

            for section, (appraisal, extractions) in zip(active_sections, deep_results):
                self.store.set_affective_state(appraisal)
                metrics.appraisal_valence = appraisal.valence
                metrics.appraisal_arousal = appraisal.arousal
                metrics.appraisal_emotion = appraisal.primary_emotion
                metrics.appraisal_intensity = appraisal.intensity
                if not extractions:
                    continue
                total_extractions += len(extractions)
                new_memories = self._create_semantic_memories(doc, encounter_id, appraisal, extractions)
                dedup_count += len(extractions) - len(new_memories)
                created_ids.extend(new_memories)
        else:
            # Larger docs: shared appraisal, parallel extraction only
            appraisal = overall_appraisal if overall_appraisal is not None else Appraisal()

            async def _parallel_extract() -> list[list[Extraction]]:
                return await asyncio.gather(*[
                    self.extractor.aextract(
                        section=s, doc=doc, appraisal=appraisal, mode=mode, max_items=max_items,
                    )
                    for s in active_sections
                ])

            section_extractions = asyncio.run(_parallel_extract())

            for section, extractions in zip(active_sections, section_extractions):
                if not extractions:
                    continue
                total_extractions += len(extractions)
                new_memories = self._create_semantic_memories(doc, encounter_id, appraisal, extractions)
                dedup_count += len(extractions) - len(new_memories)
                created_ids.extend(new_memories)

        if self.config.verbose:
            _emit(self.config, f"  Created {len(created_ids)} semantic memories")

        self.stats["files_processed"] += 1
        self.stats["memories_created"] += len(created_ids) + (1 if encounter_id else 0)

        # Store metrics
        metrics.extraction_count = total_extractions
        metrics.dedup_count = dedup_count
        metrics.memory_count = len(created_ids) + (1 if encounter_id else 0)
        metrics.llm_calls = self.llm.call_count - llm_calls_start
        metrics.duration_seconds = time.time() - metrics.start_time
        self.store.store_metrics(metrics)

        return len(created_ids)

    GIT_IGNORE_DIRS = {
        ".git", "node_modules", "__pycache__", ".venv", "venv",
        ".env", "dist", "build", ".tox", ".mypy_cache",
        ".pytest_cache", "vendor", ".bundle", ".next", "coverage",
    }

    def ingest_directory(
        self,
        dir_path: Path,
        recursive: bool = True,
        exclude_dirs: set[str] | None = None,
    ) -> int:
        if _should_cancel(self.config):
            raise RuntimeError("Ingestion cancelled")
        if not dir_path.exists() or not dir_path.is_dir():
            _emit(self.config, f"Directory not found: {dir_path}")
            return 0
        pattern = "**/*" if recursive else "*"
        files = [
            f for f in dir_path.glob(pattern)
            if f.is_file()
            and f.suffix.lower() in self.SUPPORTED_EXTENSIONS
            and (
                exclude_dirs is None
                or not any(part in exclude_dirs for part in f.relative_to(dir_path).parts)
            )
        ]
        if self.config.verbose:
            _emit(self.config, f"Found {len(files)} files to process")
        total = 0
        for file_path in files:
            total += self.ingest_file(file_path)
        return total

    def ingest_url(self, url: str, title: str | None = None) -> int:
        """Ingest content from a URL."""
        metrics = IngestionMetrics(start_time=time.time())
        llm_calls_start = self.llm.call_count

        if self.config.verbose:
            _emit(self.config, f"\nFetching: {url}")

        try:
            content = WebReader.read(url)
            metrics.source_size_bytes = len(content.encode("utf-8"))
        except Exception:
            # Fallback: try as RSS/Atom feed
            try:
                content = RssReader.read(url)
                if content:
                    metrics.source_size_bytes = len(content.encode("utf-8"))
                    if self.config.verbose:
                        _emit(self.config, "  Fetched as RSS/Atom feed")
                else:
                    raise RuntimeError("No RSS entries found")
            except Exception as exc2:
                _emit(self.config, f"  Error fetching URL: {exc2}")
                self.stats["errors"] += 1
                metrics.errors.append(str(exc2))
                return 0

        content_hash = _hash_text(content)
        words = _word_count(content)
        mode = self.config.mode

        # Extract title from content header if not provided
        if not title:
            import re
            title_match = re.search(r"\[Title: (.+?)\]", content)
            if title_match:
                title = title_match.group(1)
            else:
                title = url.split("/")[-1] or url

        metrics.word_count = words
        metrics.mode = mode.value
        metrics.source_type = "web"

        doc = DocumentInfo(
            title=title,
            source_type="web",
            content_hash=content_hash,
            word_count=words,
            path=url,
            file_type=".html",
        )

        if self.store.has_receipt(content_hash):
            if self.config.verbose:
                _emit(self.config, f"  Already ingested (hash={content_hash[:8]}...)")
            return 0

        virtual_path = Path("web_content.md")
        sections = self.sectioner.split(content, virtual_path)

        if self.config.verbose:
            _emit(self.config, f"  Mode: {mode.value} | Words: {words} | Sections: {len(sections)}")

        # FAST mode: small docs get per-section appraisal, larger get doc-level
        base_context = self._build_appraisal_context(doc)
        use_deep = words <= self.config.deep_max_words

        if use_deep:
            overall_appraisal = None
        else:
            sample = self._sample_content(content)
            overall_appraisal = self.appraiser.appraise(content=sample, context=base_context, mode=mode)
            self.store.set_affective_state(overall_appraisal)
            metrics.appraisal_valence = overall_appraisal.valence
            metrics.appraisal_arousal = overall_appraisal.arousal
            metrics.appraisal_emotion = overall_appraisal.primary_emotion
            metrics.appraisal_intensity = overall_appraisal.intensity

        encounter_id = self._create_encounter_memory(doc, overall_appraisal, mode)

        created_ids: list[str] = []
        total_extractions = 0
        dedup_count = 0

        active_sections = [s for s in sections if not self._skip_section(s.title)]
        max_items = self.config.max_facts_per_section

        if use_deep:
            async def _parallel_deep() -> list[tuple[Appraisal, list[Extraction]]]:
                async def _appraise_and_extract(s: Section) -> tuple[Appraisal, list[Extraction]]:
                    sample = self._sample_content(s.content)
                    apr = await self.appraiser.aappraise(content=sample, context=base_context, mode=mode)
                    exts = await self.extractor.aextract(
                        section=s, doc=doc, appraisal=apr, mode=mode, max_items=max_items,
                    )
                    return apr, exts
                return await asyncio.gather(*[_appraise_and_extract(s) for s in active_sections])

            deep_results = asyncio.run(_parallel_deep())

            for section, (section_appraisal, extractions) in zip(active_sections, deep_results):
                self.store.set_affective_state(section_appraisal)
                if extractions:
                    total_extractions += len(extractions)
                    new_memories = self._create_semantic_memories(doc, encounter_id, section_appraisal, extractions)
                    dedup_count += len(extractions) - len(new_memories)
                    created_ids.extend(new_memories)
        else:
            appraisal = overall_appraisal if overall_appraisal is not None else Appraisal()

            async def _parallel_extract() -> list[list[Extraction]]:
                return await asyncio.gather(*[
                    self.extractor.aextract(
                        section=s, doc=doc, appraisal=appraisal, mode=mode, max_items=max_items,
                    )
                    for s in active_sections
                ])

            section_extractions = asyncio.run(_parallel_extract())

            for section, extractions in zip(active_sections, section_extractions):
                if extractions:
                    total_extractions += len(extractions)
                    new_memories = self._create_semantic_memories(doc, encounter_id, appraisal, extractions)
                    dedup_count += len(extractions) - len(new_memories)
                    created_ids.extend(new_memories)

        if self.config.verbose:
            _emit(self.config, f"  Created {len(created_ids)} semantic memories")

        self.stats["files_processed"] += 1
        self.stats["memories_created"] += len(created_ids) + (1 if encounter_id else 0)

        metrics.extraction_count = total_extractions
        metrics.dedup_count = dedup_count
        metrics.memory_count = len(created_ids) + (1 if encounter_id else 0)
        metrics.llm_calls = self.llm.call_count - llm_calls_start
        metrics.duration_seconds = time.time() - metrics.start_time
        self.store.store_metrics(metrics)

        return len(created_ids)

    def _sample_content(self, content: str, limit: int = 2000) -> str:
        if len(content) <= limit:
            return content
        head = content[:limit]
        tail = content[-limit:]
        return f"{head}\n\n...\n\n{tail}"

    def _build_appraisal_context(self, doc: DocumentInfo) -> dict[str, Any]:
        ctx = {
            "document": {
                "title": doc.title,
                "source_type": doc.source_type,
                "word_count": doc.word_count,
            }
        }
        try:
            ctx.update(self.store.fetch_appraisal_context())
        except Exception:
            pass
        return ctx

    def _source_payload(self, doc: DocumentInfo) -> dict[str, Any]:
        now = datetime.now(timezone.utc).isoformat()
        payload = {
            "kind": doc.source_type,
            "ref": doc.content_hash,
            "label": doc.title,
            "observed_at": now,
            "content_hash": doc.content_hash,
            "path": doc.path,
        }
        if self.config.base_trust is not None:
            payload["trust"] = float(self.config.base_trust)
        return payload

    def _create_archive_encounter(self, doc: DocumentInfo) -> str | None:
        source = self._source_payload(doc)
        text = f"I have access to '{doc.title}' but haven't engaged with it yet."
        context = {
            "activity": "archived",
            "source_type": doc.source_type,
            "source_ref": doc.content_hash,
            "word_count": doc.word_count,
            "mode": "archived",
            "awaiting_processing": True,
        }
        importance = max(self.config.min_importance_floor or 0.0, 0.2)
        encounter_id = self.store.create_encounter_memory(
            text=text,
            source=source,
            emotional_valence=0.0,
            context=context,
            importance=importance,
        )
        self._apply_decay(encounter_id, intensity=0.0)
        return encounter_id

    def _create_encounter_memory(self, doc: DocumentInfo, appraisal: Appraisal | None, mode: IngestionMode) -> str | None:
        source = self._source_payload(doc)
        appraisal = appraisal or Appraisal()
        summary = appraisal.summary or ""
        if not summary:
            summary = f"It felt {appraisal.primary_emotion} with intensity {appraisal.intensity:.2f}."
        text = f"I read '{doc.title}'. {summary}"
        context = {
            "activity": "reading",
            "source_type": doc.source_type,
            "source_ref": doc.content_hash,
            "word_count": doc.word_count,
            "mode": mode.value,
            "appraisal": appraisal.__dict__,
        }
        importance = max(self.config.min_importance_floor or 0.0, 0.4 + appraisal.intensity * 0.4)
        encounter_id = self.store.create_encounter_memory(
            text=text,
            source=source,
            emotional_valence=appraisal.valence,
            context=context,
            importance=importance,
        )
        self._apply_decay(encounter_id, intensity=appraisal.intensity)
        return encounter_id

    def _apply_decay(self, memory_id: str, intensity: float) -> None:
        if self.config.permanent:
            self.store.update_decay_rate(memory_id, 0.0)
            return
        decay = _decay_rate_for_intensity(intensity)
        self.store.update_decay_rate(memory_id, decay)

    def _create_semantic_memories(
        self,
        doc: DocumentInfo,
        encounter_id: str | None,
        appraisal: Appraisal,
        extractions: list[Extraction],
    ) -> list[str]:
        created: list[str] = []
        source = self._source_payload(doc)

        # Pre-warm embedding cache: batch all extraction texts into a single
        # call so subsequent per-extraction recall + create are cache hits.
        texts_to_embed = [
            ext.content
            for ext in extractions
            if ext.confidence >= self.config.min_confidence_threshold
        ]
        if texts_to_embed:
            self.store.prefetch_embeddings(texts_to_embed)

        # Collect deferred work for batch execution after the per-item loop
        concept_pairs: list[tuple[str, str]] = []
        worldview_hints: dict[str, str | list] = {}  # hint_key -> raw hint
        deferred_worldview_edges: list[tuple[str, str | list, str, float]] = []  # (memory_id, hint, rel_type_name, confidence)
        deferred_edges: list[tuple[str, str, RelationshipType, float]] = []

        # Dedup/related/create routing is DB-owned (db/41 ingest_route_extractions):
        # thresholds live in config and the nearest-neighbor search runs in one
        # batched call instead of a Python recall per extraction.
        plan = self.store.route_extractions(extractions, self.config.min_confidence_threshold)
        plan_by_index = {p["index"]: p for p in plan if isinstance(p, dict) and "index" in p}

        for idx, ext in enumerate(extractions):
            routed = plan_by_index.get(idx)
            if routed is None:
                continue  # dropped below the confidence threshold by the router
            decision = routed.get("decision")
            matched_id = routed.get("matched_memory_id")

            if decision == "duplicate" and matched_id:
                # Corroboration, not re-creation: merge the source, revise
                # confidence via the audited policy, and keep an evidence edge
                # from the encounter memory (#34/#35).
                try:
                    self.store.add_evidence(
                        str(matched_id),
                        "supports",
                        source,
                        evidence_memory_id=encounter_id,
                        context="fast_ingest",
                    )
                except Exception:
                    logger.exception("fast ingest corroboration failed for memory %s", matched_id)
                continue

            importance = ext.importance
            if self.config.min_importance_floor is not None:
                importance = max(importance, self.config.min_importance_floor)
            trust = self.config.base_trust

            memory_id = self.store.create_semantic_memory(
                content=ext.content,
                confidence=ext.confidence,
                category=ext.category,
                related_concepts=ext.connections,
                source=source,
                importance=importance,
                trust=trust,
            )
            created.append(memory_id)

            # Collect concept links for batch
            for concept in ext.concepts:
                concept_pairs.append((memory_id, concept.strip()))

            # Collect worldview edge hints (deduplicated lookup later)
            if ext.supports:
                hint_key = str(ext.supports) if isinstance(ext.supports, list) else ext.supports
                worldview_hints[hint_key] = ext.supports
                deferred_worldview_edges.append((memory_id, ext.supports, "SUPPORTS", ext.confidence))

            if ext.contradicts:
                hint_key = str(ext.contradicts) if isinstance(ext.contradicts, list) else ext.contradicts
                worldview_hints[hint_key] = ext.contradicts
                deferred_worldview_edges.append((memory_id, ext.contradicts, "CONTRADICTS", ext.confidence))

            if encounter_id:
                deferred_edges.append((memory_id, encounter_id, RelationshipType.DERIVED_FROM, 0.9))
            if decision == "related" and matched_id:
                deferred_edges.append((memory_id, str(matched_id), RelationshipType.ASSOCIATED, 0.6))
            self._apply_decay(memory_id, intensity=appraisal.intensity)

        # --- Batch flush phase ---

        # 1. Batch concept linking
        if concept_pairs:
            try:
                self.store.link_concepts_batch(concept_pairs)
            except Exception:
                pass

        # 2. Batch worldview lookups (deduplicated)
        worldview_cache: dict[str, str | None] = {}
        for hint_key, hint_val in worldview_hints.items():
            if hint_key not in worldview_cache:
                worldview_cache[hint_key] = self._find_worldview_by_content(hint_val)

        # 3. Resolve worldview edges and add to deferred_edges
        for memory_id, hint, rel_type_name, confidence in deferred_worldview_edges:
            hint_key = str(hint) if isinstance(hint, list) else hint
            worldview_id = worldview_cache.get(hint_key)
            if worldview_id:
                rel_type = RelationshipType.SUPPORTS if rel_type_name == "SUPPORTS" else RelationshipType.CONTRADICTS
                deferred_edges.append((memory_id, worldview_id, rel_type, confidence))

        # 4. Batch relationship creation
        if deferred_edges:
            try:
                self.store.connect_memories_batch(deferred_edges)
            except Exception:
                pass

        return created

    def _find_worldview_by_content(self, hint) -> str | None:
        """Find a worldview memory matching the given hint."""
        if isinstance(hint, list):
            hint = " ".join(str(h) for h in hint if h) if hint else ""
        if not hint or not isinstance(hint, str) or not hint.strip():
            return None
        try:
            results = self.store.client.recall(
                hint.strip(),
                limit=3,
                memory_types=[ApiMemoryType.WORLDVIEW],
            )
            for mem in results.memories:
                if mem.similarity is not None and mem.similarity >= 0.7:
                    return str(mem.id)
        except Exception:
            pass
        return None

    def _skip_section(self, title: str) -> bool:
        lowered = title.strip().lower()
        return any(skip in lowered for skip in self.config.skip_sections)

    def _run_rlm_ingest(
        self,
        mode: IngestionMode,
        doc: DocumentInfo,
        sections: list[Section],
        metrics: IngestionMetrics,
        llm_calls_start: int,
    ) -> int:
        """Run slow or hybrid ingestion via RLM loop.

        This is a sync wrapper that calls the async slow_ingest_rlm functions.
        Used by ingest_file() when mode is SLOW or HYBRID.
        """
        import asyncio as _asyncio
        import concurrent.futures

        from services.slow_ingest_rlm import run_hybrid_ingest, run_slow_ingest

        if self.config.dsn:
            dsn = self.config.dsn
        else:
            dsn = (
                f"postgresql://{self.config.db_user}:{self.config.db_password}"
                f"@{self.config.db_host}:{self.config.db_port}/{self.config.db_name}"
            )
        llm_cfg = self.llm._cfg

        runner = run_slow_ingest if mode == IngestionMode.SLOW else run_hybrid_ingest
        coro = runner(
            pipeline=self,
            doc=doc,
            sections=sections,
            llm_config=llm_cfg,
            dsn=dsn,
        )

        # Run async coroutine from sync context. If an event loop is already
        # running (e.g. called from ingest_api thread), offload to a fresh
        # thread with its own loop; otherwise use asyncio.run() directly.
        try:
            _asyncio.get_running_loop()
            # Already inside an event loop -- run in a separate thread.
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
                result = pool.submit(_asyncio.run, coro).result()
        except RuntimeError:
            # No running loop -- safe to use asyncio.run().
            result = _asyncio.run(coro)

        count = result.get("memories_created", 0)
        self.stats["files_processed"] += 1
        self.stats["memories_created"] += count

        metrics.memory_count = count
        metrics.mode = mode.value
        metrics.llm_calls = self.llm.call_count - llm_calls_start
        metrics.duration_seconds = time.time() - metrics.start_time
        self.store.store_metrics(metrics)

        if self.config.verbose:
            _emit(self.config, f"  RLM {mode.value} ingest: {count} memories created")
            if mode == IngestionMode.HYBRID:
                _emit(
                    self.config,
                    f"  Slow chunks: {result.get('slow_chunks', 0)} | "
                    f"Fast chunks: {result.get('fast_chunks', 0)}",
                )

        return count

    def check_and_process_archived(self, query: str, threshold: float = 0.75) -> list[str]:
        """
        Check if any archived content matches the query and process it.

        This implements retrieval-triggered processing: when a query surfaces
        archived content that hasn't been fully processed, we upgrade it now.

        Returns list of content hashes that were processed.
        """
        if self.store.client is None:
            self.store.connect()

        # Find archived content matching the query
        rows = self.store._fetchval(
            """
            SELECT jsonb_agg(jsonb_build_object(
                'memory_id', memory_id,
                'content_hash', content_hash,
                'title', title,
                'similarity', similarity,
                'source_path', source_path
            ))
            FROM check_archived_for_query($1, $2, 5)
            """,
            query,
            threshold,
        )

        if not rows:
            return []

        archived = json.loads(rows) if isinstance(rows, str) else rows
        if not archived:
            return []

        processed_hashes: list[str] = []

        for item in archived:
            if not item:
                continue

            content_hash = item.get("content_hash")
            source_path = item.get("source_path")
            title = item.get("title")
            memory_id = item.get("memory_id")

            if not content_hash:
                continue

            _emit(self.config, f"Processing archived content triggered by query: {title}")

            # Attempt to re-read the source file if it exists
            if source_path and source_path != "stdin" and not source_path.startswith("http"):
                path = Path(source_path)
                if path.exists():
                    # Re-ingest the file with the current mode (not archive)
                    original_mode = self.config.mode
                    self.config.mode = IngestionMode.FAST
                    try:
                        # Mark as processed first to avoid duplicate detection
                        self.store._fetchval(
                            "SELECT mark_archived_as_processed($1::uuid)",
                            memory_id,
                        )
                        self.ingest_file(path)
                        processed_hashes.append(content_hash)
                    finally:
                        self.config.mode = original_mode
                    continue

            # If source file not available, just mark as processed
            self.store._fetchval(
                "SELECT mark_archived_as_processed($1::uuid)",
                memory_id,
            )
            processed_hashes.append(content_hash)

        return processed_hashes

    def print_stats(self) -> None:
        _emit(self.config, "\n" + "=" * 50)
        _emit(self.config, "INGESTION COMPLETE")
        _emit(self.config, "=" * 50)
        _emit(self.config, f"Files processed:   {self.stats['files_processed']}")
        _emit(self.config, f"Memories created:  {self.stats['memories_created']}")
        _emit(self.config, f"Errors:            {self.stats['errors']}")
        _emit(self.config, "=" * 50)

    def close(self) -> None:
        self.store.close()


# =========================================================================
# ARCHIVED CONTENT PROCESSOR
# =========================================================================


class ArchivedContentProcessor:
    """
    Processor for upgrading archived content to full memories.

    This can be used:
    1. During recall - when a query surfaces relevant archived content
    2. By maintenance workers - batch processing of pending archives
    3. By CLI - manual processing of specific content
    """

    def __init__(self, config: Config):
        self.config = config
        self.pipeline = IngestionPipeline(config)

    def process_for_query(self, query: str, threshold: float = 0.75) -> list[str]:
        """
        Check if archived content matches a query and process it.

        Returns list of content hashes that were processed.
        """
        return self.pipeline.check_and_process_archived(query, threshold)

    def process_by_hash(self, content_hash: str) -> bool:
        """Process a specific archived item by content hash."""
        archived = self.pipeline.store.check_archived_for_query(
            content_hash, threshold=0.0, limit=1
        )

        if not archived:
            # Try direct lookup
            row = self.pipeline.store._fetchval(
                """
                SELECT jsonb_build_object(
                    'memory_id', id,
                    'content_hash', source_attribution->>'content_hash',
                    'title', source_attribution->>'label',
                    'source_path', source_attribution->>'path'
                )
                FROM memories
                WHERE type = 'episodic'
                  AND source_attribution->>'content_hash' = $1
                  AND metadata->>'awaiting_processing' = 'true'
                LIMIT 1
                """,
                content_hash,
            )
            if not row:
                return False
            archived = [json.loads(row) if isinstance(row, str) else row]

        for item in archived:
            if not item:
                continue

            source_path = item.get("source_path")
            memory_id = item.get("memory_id")

            if source_path and source_path != "stdin" and not source_path.startswith("http"):
                path = Path(source_path)
                if path.exists():
                    self.pipeline.store.mark_archived_processed(memory_id)
                    original_mode = self.config.mode
                    self.config.mode = IngestionMode.FAST
                    try:
                        self.pipeline.ingest_file(path)
                    finally:
                        self.config.mode = original_mode
                    return True

            # Mark as processed even if file not found
            return self.pipeline.store.mark_archived_processed(memory_id)

        return False

    def process_batch(self, limit: int = 10) -> int:
        """Process a batch of archived items."""
        rows = self.pipeline.store._fetchval(
            """
            SELECT ARRAY_AGG(source_attribution->>'content_hash')
            FROM memories
            WHERE type = 'episodic'
              AND metadata->>'awaiting_processing' = 'true'
            ORDER BY importance DESC, created_at ASC
            LIMIT $1
            """,
            limit,
        )

        if not rows:
            return 0

        hashes = list(rows) if rows else []
        count = 0
        for h in hashes:
            if h and self.process_by_hash(h):
                count += 1

        return count

    def close(self) -> None:
        self.pipeline.close()


# =========================================================================
# CLI
# =========================================================================


def _get_db_env_defaults() -> dict[str, Any]:
    """Get database configuration from environment variables."""
    env_db_port_raw = os.getenv("POSTGRES_PORT")
    try:
        env_db_port = int(env_db_port_raw) if env_db_port_raw else 43815
    except ValueError:
        env_db_port = 43815
    return {
        "db_host": os.getenv("POSTGRES_HOST", "localhost"),
        "db_port": env_db_port,
        "db_name": os.getenv("POSTGRES_DB", "hexis_memory"),
        "db_user": os.getenv("POSTGRES_USER", "postgres"),
        "db_password": os.getenv("POSTGRES_PASSWORD", "password"),
    }


def _add_common_args(parser: argparse.ArgumentParser, env_defaults: dict[str, Any]) -> None:
    """Add common arguments shared across subcommands."""
    parser.add_argument("--endpoint", "-e", default=None, help="LLM endpoint (overrides DB config)")
    parser.add_argument("--model", "-m", default=None, help="LLM model name (overrides DB config)")
    parser.add_argument("--api-key", default=None, help="LLM API key (overrides DB config)")
    parser.add_argument("--provider", default=None, help="LLM provider (overrides DB config)")

    parser.add_argument("--db-host", default=env_defaults["db_host"], help="Database host")
    parser.add_argument("--db-port", type=int, default=env_defaults["db_port"], help="Database port")
    parser.add_argument("--db-name", default=env_defaults["db_name"], help="Database name")
    parser.add_argument("--db-user", default=env_defaults["db_user"], help="Database user")
    parser.add_argument("--db-password", default=env_defaults["db_password"], help="Database password")

    parser.add_argument("--quiet", "-q", action="store_true", help="Suppress verbose output")


def _load_llm_config_from_db(args: argparse.Namespace) -> dict[str, Any] | None:
    """Load fully-resolved LLM config using the same path as `hexis chat`.

    Uses core.llm_config.load_llm_config() which handles provider-specific
    credential loading (OAuth tokens for Codex, etc.).
    """
    import asyncio

    import asyncpg

    async def _load() -> dict[str, Any] | None:
        try:
            dsn = (
                f"postgresql://{args.db_user}:{args.db_password}"
                f"@{args.db_host}:{args.db_port}/{args.db_name}"
            )
            conn = await asyncpg.connect(dsn)
            try:
                from core.llm_config import load_llm_config

                cfg = await load_llm_config(conn, "llm.chat", fallback_key="llm")
                return cfg
            finally:
                await conn.close()
        except Exception:
            return None

    return asyncio.run(_load())


def _build_config_from_args(args: argparse.Namespace) -> Config:
    """Build Config from parsed arguments.

    Priority: CLI flags > DB config (with full credential resolution) > defaults.
    """
    # Load fully-resolved config from DB (handles OAuth, etc.)
    db_llm = _load_llm_config_from_db(args) or {}

    # CLI overrides (only if explicitly set — None means "not provided")
    cli_endpoint = getattr(args, "endpoint", None)
    cli_model = getattr(args, "model", None)
    cli_api_key = getattr(args, "api_key", None)
    cli_provider = getattr(args, "provider", None)

    llm_config: dict[str, Any] = dict(db_llm)
    if cli_endpoint:
        llm_config["endpoint"] = cli_endpoint
    if cli_model:
        llm_config["model"] = cli_model
    if cli_api_key:
        llm_config["api_key"] = cli_api_key
    if cli_provider:
        llm_config["provider"] = cli_provider

    return Config(
        llm_config=llm_config,
        db_host=args.db_host,
        db_port=args.db_port,
        db_name=args.db_name,
        db_user=args.db_user,
        db_password=args.db_password,
        mode=_normalize_mode(getattr(args, "mode", "auto")),
        min_importance_floor=getattr(args, "min_importance", None),
        permanent=getattr(args, "permanent", False),
        base_trust=getattr(args, "base_trust", None),
        verbose=not getattr(args, "quiet", False),
    )


def _cmd_ingest(args: argparse.Namespace) -> None:
    """Handle the ingest subcommand."""
    config = _build_config_from_args(args)
    pipeline = IngestionPipeline(config)

    try:
        if args.stdin:
            count = _ingest_stdin(pipeline, args)
        elif args.url:
            count = _ingest_url(pipeline, args)
        elif getattr(args, "github", None):
            count = _ingest_github(pipeline, args)
        elif args.file:
            count = pipeline.ingest_file(args.file)
        elif args.input:
            count = pipeline.ingest_directory(args.input, recursive=not args.no_recursive)
        else:
            print("Error: No input source specified")
            return
        pipeline.print_stats()
    except KeyboardInterrupt:
        print("\nInterrupted by user")
    finally:
        pipeline.close()


def _ingest_stdin(pipeline: IngestionPipeline, args: argparse.Namespace) -> int:
    """Ingest content from stdin."""
    content = sys.stdin.read()
    if not content.strip():
        _emit(pipeline.config, "No content received from stdin")
        return 0

    content_type = getattr(args, "stdin_type", "text") or "text"
    title = getattr(args, "stdin_title", None) or f"stdin-{_hash_text(content)[:8]}"

    # Create a virtual DocumentInfo
    content_hash = _hash_text(content)
    words = _word_count(content)
    mode = pipeline.config.mode

    source_type_map = {
        "text": "document",
        "markdown": "document",
        "code": "code",
        "json": "data",
        "yaml": "data",
        "data": "data",
    }
    source_type = source_type_map.get(content_type, "document")

    doc = DocumentInfo(
        title=title,
        source_type=source_type,
        content_hash=content_hash,
        word_count=words,
        path="stdin",
        file_type=f".{content_type}",
    )

    if pipeline.store.has_receipt(content_hash):
        _emit(pipeline.config, f"Content already ingested (hash={content_hash[:8]}...)")
        return 0

    # Create virtual path for sectioning
    virtual_path = Path(f"stdin.{content_type}")
    sections = pipeline.sectioner.split(content, virtual_path)

    _emit(pipeline.config, f"Processing stdin: {title}")
    _emit(pipeline.config, f"  Mode: {mode.value} | Words: {words} | Sections: {len(sections)}")

    # FAST mode: small docs get per-section appraisal, larger get doc-level
    base_context = pipeline._build_appraisal_context(doc)
    use_deep = words <= pipeline.config.deep_max_words

    if use_deep:
        overall_appraisal = None
    else:
        sample = pipeline._sample_content(content)
        overall_appraisal = pipeline.appraiser.appraise(content=sample, context=base_context, mode=mode)
        pipeline.store.set_affective_state(overall_appraisal)

    encounter_id = pipeline._create_encounter_memory(doc, overall_appraisal, mode)

    created_ids: list[str] = []
    for section in sections:
        if pipeline._skip_section(section.title):
            continue
        if use_deep:
            sample = pipeline._sample_content(section.content)
            section_appraisal = pipeline.appraiser.appraise(content=sample, context=base_context, mode=mode)
            pipeline.store.set_affective_state(section_appraisal)
        else:
            section_appraisal = overall_appraisal if overall_appraisal is not None else Appraisal()

        max_items = pipeline.config.max_facts_per_section

        extractions = pipeline.extractor.extract(
            section=section,
            doc=doc,
            appraisal=section_appraisal,
            mode=mode,
            max_items=max_items,
        )
        if extractions:
            created_ids.extend(pipeline._create_semantic_memories(doc, encounter_id, section_appraisal, extractions))

    _emit(pipeline.config, f"  Created {len(created_ids)} semantic memories")
    pipeline.stats["files_processed"] += 1
    pipeline.stats["memories_created"] += len(created_ids) + (1 if encounter_id else 0)
    return len(created_ids)


def _ingest_url(pipeline: IngestionPipeline, args: argparse.Namespace) -> int:
    """Ingest content from a URL."""
    title = getattr(args, "title", None)
    return pipeline.ingest_url(args.url, title=title)


def _ingest_github(pipeline: IngestionPipeline, args: argparse.Namespace) -> int:
    """Clone a GitHub repo to a temp dir and ingest its contents."""
    import shutil
    import subprocess
    import tempfile

    repo = args.github.strip()
    # Accept owner/repo shorthand
    if "/" in repo and not repo.startswith("http"):
        repo = f"https://github.com/{repo}"

    tmpdir = tempfile.mkdtemp(prefix="hexis_git_")
    try:
        cmd = ["git", "clone", "--depth", "1"]
        branch = getattr(args, "branch", None)
        if branch:
            cmd.extend(["--branch", branch])
        cmd.extend([repo, tmpdir])

        _emit(pipeline.config, f"Cloning {repo} ...")
        subprocess.check_call(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.PIPE)

        count = pipeline.ingest_directory(
            Path(tmpdir),
            recursive=True,
            exclude_dirs=IngestionPipeline.GIT_IGNORE_DIRS,
        )
        _emit(pipeline.config, f"Ingested {count} memories from {repo}")
        return count
    except subprocess.CalledProcessError as e:
        _emit(pipeline.config, f"Git clone failed: {e}")
        return 0
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)


def _cmd_status(args: argparse.Namespace) -> None:
    """Handle the status subcommand."""
    config = _build_config_from_args(args)
    store = MemoryStore(config)
    store.connect()

    try:
        if args.pending:
            # Query for archived/pending memories
            rows = store._fetchval(
                """
                SELECT jsonb_agg(jsonb_build_object(
                    'id', id,
                    'title', source_attribution->>'label',
                    'hash', source_attribution->>'content_hash',
                    'created_at', created_at
                ))
                FROM memories
                WHERE type = 'episodic'
                  AND metadata->>'awaiting_processing' = 'true'
                ORDER BY created_at DESC
                LIMIT 50
                """
            )
            pending = json.loads(rows) if rows else []

            if args.json:
                print(json.dumps(pending, indent=2, default=str))
            else:
                if not pending:
                    print("No pending ingestions")
                else:
                    print(f"Pending ingestions: {len(pending)}")
                    for p in pending:
                        print(f"  - {p.get('title', 'Unknown')} ({p.get('hash', '')[:8]}...)")
        else:
            # General ingestion stats
            stats = store._fetchval(
                """
                SELECT jsonb_build_object(
                    'total_memories', (SELECT COUNT(*) FROM memories),
                    'episodic', (SELECT COUNT(*) FROM memories WHERE type = 'episodic'),
                    'semantic', (SELECT COUNT(*) FROM memories WHERE type = 'semantic'),
                    'pending', (SELECT COUNT(*) FROM memories WHERE type = 'episodic' AND metadata->>'awaiting_processing' = 'true'),
                    'recent_24h', (SELECT COUNT(*) FROM memories WHERE created_at > NOW() - INTERVAL '24 hours')
                )
                """
            )
            stats_data = json.loads(stats) if stats else {}

            if args.json:
                print(json.dumps(stats_data, indent=2))
            else:
                print("Ingestion Status:")
                print(f"  Total memories:     {stats_data.get('total_memories', 0)}")
                print(f"  Episodic memories:  {stats_data.get('episodic', 0)}")
                print(f"  Semantic memories:  {stats_data.get('semantic', 0)}")
                print(f"  Pending processing: {stats_data.get('pending', 0)}")
                print(f"  Last 24 hours:      {stats_data.get('recent_24h', 0)}")
    finally:
        store.close()


def _cmd_process(args: argparse.Namespace) -> None:
    """Handle the process subcommand - upgrade archived content."""
    config = _build_config_from_args(args)
    processor = ArchivedContentProcessor(config)

    try:
        if args.content_hash:
            success = processor.process_by_hash(args.content_hash)
            print(f"Processed: {'Yes' if success else 'No (not found or failed)'}")
        elif args.all_archived:
            count = processor.process_batch(limit=getattr(args, "limit", 10))
            print(f"Processed {count} archived items")
        else:
            print("Error: Specify --content-hash or --all-archived")
    except KeyboardInterrupt:
        print("\nInterrupted by user")
    finally:
        processor.close()


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Hexis Universal Ingestion Pipeline",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Subcommands:
  ingest   Ingest files, directories, URLs, or stdin
  status   Show ingestion status and pending items
  process  Process archived content that hasn't been fully engaged

Examples:
  %(prog)s ingest --file doc.md --mode fast
  %(prog)s ingest --input ./docs --mode slow
  %(prog)s ingest --url https://example.com/article
  echo "Some text" | %(prog)s ingest --stdin --stdin-type text
  %(prog)s status --pending
  %(prog)s process --all-archived
        """,
    )

    env_defaults = _get_db_env_defaults()
    subparsers = parser.add_subparsers(dest="subcommand")

    # Ingest subcommand
    ingest_p = subparsers.add_parser("ingest", help="Ingest content into memory")
    input_group = ingest_p.add_mutually_exclusive_group()
    input_group.add_argument("--file", "-f", type=Path, help="Single file to ingest")
    input_group.add_argument("--input", "-i", type=Path, help="Directory to ingest")
    input_group.add_argument("--url", "-u", type=str, help="URL to fetch and ingest")
    input_group.add_argument("--stdin", action="store_true", help="Read content from stdin")
    input_group.add_argument("--github", "-g", type=str, help="GitHub repo URL or owner/repo to clone and ingest")

    ingest_p.add_argument("--branch", type=str, default=None, help="Branch to clone (default: repo default)")
    ingest_p.add_argument("--stdin-type", choices=["text", "markdown", "code", "json", "yaml", "data"], default="text", help="Content type for stdin input")
    ingest_p.add_argument("--stdin-title", type=str, help="Title for stdin content")
    ingest_p.add_argument("--title", type=str, help="Override document title")

    ingest_p.add_argument("--mode", default="fast", choices=[m.value for m in IngestionMode], help="Ingestion mode")
    ingest_p.add_argument("--no-recursive", action="store_true", help="Don't recurse into subdirectories")
    ingest_p.add_argument("--min-importance", type=float, help="Minimum importance floor")
    ingest_p.add_argument("--permanent", action="store_true", help="Mark memories as permanent (no decay)")
    ingest_p.add_argument("--base-trust", type=float, help="Base trust level for source")

    _add_common_args(ingest_p, env_defaults)

    # Status subcommand
    status_p = subparsers.add_parser("status", help="Show ingestion status")
    status_p.add_argument("--pending", action="store_true", help="Show pending/archived ingestions")
    status_p.add_argument("--json", action="store_true", help="Output as JSON")
    _add_common_args(status_p, env_defaults)

    # Process subcommand
    process_p = subparsers.add_parser("process", help="Process archived content")
    process_p.add_argument("--content-hash", type=str, help="Content hash of specific archived item")
    process_p.add_argument("--all-archived", action="store_true", help="Process all archived items")
    process_p.add_argument("--limit", type=int, default=10, help="Max items to process")
    _add_common_args(process_p, env_defaults)

    args = parser.parse_args()

    # Default to ingest if no subcommand (for backwards compatibility)
    if args.subcommand is None:
        # Check if any ingest-style args were provided
        if hasattr(args, "file") or hasattr(args, "input"):
            parser.print_help()
            print("\nError: Please use 'ingest' subcommand. Example: python -m services.ingest ingest --file doc.md")
        else:
            parser.print_help()
        return

    if args.subcommand == "ingest":
        if not (args.file or args.input or args.url or args.stdin or getattr(args, "github", None)):
            print("Error: One of --file, --input, --url, --stdin, or --github is required")
            return
        _cmd_ingest(args)
    elif args.subcommand == "status":
        _cmd_status(args)
    elif args.subcommand == "process":
        _cmd_process(args)


if __name__ == "__main__":
    main()
